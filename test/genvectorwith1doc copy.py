#!/usr/bin/env python3
# -*- coding: utf-8 -*-

"""
Single Document Vectorizer
--------------------------
A step-by-step tool to process a single document, extract text, create chunks, 
and generate vector embeddings for retrieval-augmented generation (RAG) systems.
Based on the batch Vector Store Generator framework.
"""

import os
import sys
import time
import json
import logging
import re
import traceback
from pathlib import Path
from typing import List, Dict, Any, Optional
from datetime import datetime

# Document processing
import docx
import PyPDF2
import pandas as pd
try:
    import pptx
except ImportError:
    pptx = None
try:
    import openpyxl
except ImportError:
    openpyxl = None

# ML/Vector libraries
from langchain_huggingface import HuggingFaceEmbeddings
from langchain_community.vectorstores import FAISS
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.schema import Document

# Terminal UI
from rich.console import Console
from rich.table import Table
from rich.panel import Panel
from rich.progress import Progress, SpinnerColumn, TextColumn, BarColumn
from rich.prompt import Prompt, Confirm
from rich import print as rprint
from rich.syntax import Syntax

# NLP support
try:
    import spacy
except ImportError:
    spacy = None

# Initialize console
console = Console()

# Configure logging
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(levelname)s - %(message)s',
    handlers=[
        logging.FileHandler('doc_vectorizer.log', encoding='utf-8'),
        logging.StreamHandler()
    ]
)
logger = logging.getLogger("Doc_Vectorizer")


class Config:
    """Configuration for document vectorizer."""
    
    def __init__(self):
        # Paths
        self.source_file = os.path.join(os.getcwd(), "sample.pdf")
        self.vector_db_path = os.path.join(os.getcwd(), "vector_single")
        self.plain_text_path = os.path.join(os.getcwd(), "plain_single")
        self.chunk_path = os.path.join(os.getcwd(), "chunk_single")
        
        # Chunking settings
        self.chunk_size = 1000
        self.chunk_overlap = 200
        
        # Language settings
        self.languages = ["en", "ko"]
        self.language_models = {}
        
        # Embedding model settings
        self.embedding_model = "all-MiniLM-L6-v2"
        
        # File info
        self.current_file_info = {
            'hash': None,
            'language': 'en',
            'processed_date': None,
            'chunks': 0,
            'plain_text_path': None,
            'chunks_directory': None
        }
        
    def save(self, filename="doc_vectorizer_config.json"):
        """Save configuration to a file."""
        # Extract only serializable data
        save_data = {k: v for k, v in self.__dict__.items() 
                    if not k.startswith('_') and k != 'language_models'}
        
        try:
            with open(filename, 'w', encoding='utf-8') as f:
                json.dump(save_data, f, indent=4)
            logger.info(f"Configuration saved to {filename}")
            return True
        except Exception as e:
            logger.error(f"Failed to save configuration: {e}")
            return False
    
    def load(self, filename="doc_vectorizer_config.json"):
        """Load configuration from a file."""
        if os.path.exists(filename):
            try:
                with open(filename, 'r', encoding='utf-8') as f:
                    data = json.load(f)
                    for key, value in data.items():
                        if hasattr(self, key):
                            setattr(self, key, value)
                logger.info(f"Configuration loaded from {filename}")
                return True
            except Exception as e:
                logger.error(f"Failed to load configuration: {e}")
                return False
        else:
            logger.warning(f"Configuration file {filename} not found")
            return False


class DocumentProcessor:
    """Process different document types and extract text content."""
    
    @staticmethod
    def extract_text_from_pdf_orig(file_path: Path) -> str:
        """Extract text from PDF"""
        text = ""
        try:
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                for page in pdf_reader.pages:
                    page_text = page.extract_text()
                    if page_text:
                        text += page_text + "\n"
        except Exception as e:
            logger.error(f"PDF extraction error for {file_path.name}: {e}")
        return text
        
    @staticmethod
    def extract_text_from_pdf_with_ocr(file_path: Path) -> str:
        """Extract text from PDF using OCR if needed"""
        # Try normal extraction first
        text = DocumentProcessor.extract_text_from_pdf(file_path)
        
        # If no text was extracted, try OCR
        if not text or not text.strip():
            try:
                import pytesseract
                from pdf2image import convert_from_path
                
                console.print("[yellow]Standard text extraction failed. Attempting OCR...[/yellow]")
                
                # Convert PDF to images
                images = convert_from_path(file_path)
                
                # Use OCR to extract text from each page
                text = ""
                for i, image in enumerate(images):
                    # For Korean documents, specify the language
                    page_text = pytesseract.image_to_string(image, lang='kor+eng')
                    text += page_text + "\n\n"
                
                if text.strip():
                    console.print("[green]OCR extraction successful![/green]")
                else:
                    console.print("[red]OCR extraction failed to find text[/red]")
                    
            except ImportError:
                console.print("[red]OCR libraries not installed. Install with:[/red]")
                console.print("pip install pytesseract pdf2image")
                console.print("[yellow]Note: You also need to install Tesseract OCR on your system[/yellow]")
                console.print("https://github.com/UB-Mannheim/tesseract/wiki")
            except Exception as e:
                console.print(f"[red]OCR extraction error: {str(e)}[/red]")
        
        return text

    @staticmethod
    def extract_text_from_pdf(file_path: Path) -> str:
        """Extract text from PDF without OCR"""
        text = ""
        try:
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                for page in pdf_reader.pages:
                    page_text = page.extract_text()
                    if page_text:
                        text += page_text + "\n"
        except Exception as e:
            logger.error(f"PDF extraction error for {file_path.name}: {e}")
        return text

    @staticmethod
    def extract_text_from_pdf_with_fallback(file_path: Path) -> str:
        """Extract text from PDF with OCR fallback"""
        # First try regular extraction
        text = DocumentProcessor.extract_text_from_pdf(file_path)
        
        # If no text was found, try OCR
        if not text or not text.strip():
            try:
                import pytesseract
                from pdf2image import convert_from_path
                
                console.print("[yellow]Standard text extraction failed. Attempting OCR...[/yellow]")
                
                try:
                    # Convert PDF to images
                    images = convert_from_path(file_path)
                    
                    # Use OCR to extract text from each page
                    text = ""
                    for i, image in enumerate(images):
                        # For Korean documents, specify the language
                        page_text = pytesseract.image_to_string(image, lang='kor+eng')
                        text += page_text + "\n\n"
                    
                    if text.strip():
                        console.print("[green]OCR extraction successful![/green]")
                    else:
                        console.print("[red]OCR extraction failed to find text[/red]")
                
                except Exception as e:
                    if "poppler" in str(e).lower():
                        console.print("[red]Poppler is not installed or not in PATH[/red]")
                        console.print("[yellow]Poppler installation instructions:[/yellow]")
                        console.print("- Windows: Download from https://github.com/oschwartz10612/poppler-windows/releases/")
                        console.print("  Extract and add 'bin' folder to PATH")
                        console.print("- macOS: Run 'brew install poppler'")
                        console.print("- Linux: Run 'sudo apt-get install poppler-utils'")
                    else:
                        console.print(f"[red]OCR extraction error: {str(e)}[/red]")
                    logger.error(f"OCR extraction error: {e}")
                    
            except ImportError:
                console.print("[red]OCR libraries not installed. Install with:[/red]")
                console.print("pip install pytesseract pdf2image")
                console.print("[yellow]Note: You also need to install Tesseract OCR and Poppler on your system[/yellow]")
                console.print("- Tesseract: https://github.com/UB-Mannheim/tesseract/wiki")
                console.print("- Poppler: https://github.com/oschwartz10612/poppler-windows/releases/ (Windows)")
                console.print("For Korean documents, download the Korean language data pack")
        
        return text

    @staticmethod
    def extract_text_from_docx(file_path: Path) -> str:
        """Extract text from DOCX"""
        text = ""
        try:
            doc = docx.Document(str(file_path))
            for paragraph in doc.paragraphs:
                text += paragraph.text + "\n"
        except Exception as e:
            logger.error(f"DOCX extraction error for {file_path.name}: {e}")
        return text
    
    @staticmethod
    def extract_text_from_txt(file_path: Path) -> str:
        """Extract text from TXT with encoding detection"""
        encodings = ['utf-8', 'cp949', 'cp1252', 'latin1', 'euc-kr']
        for encoding in encodings:
            try:
                with open(file_path, 'r', encoding=encoding) as f:
                    text = f.read()
                    if text.strip():
                        return text
            except UnicodeDecodeError:
                continue
        return ""
    
    @staticmethod
    def extract_text_from_csv(file_path: Path) -> str:
        """Extract text from CSV with multiple encoding attempts"""
        encodings = ['utf-8', 'cp949', 'cp1252', 'latin1', 'euc-kr']
        
        for encoding in encodings:
            try:
                df = pd.read_csv(file_path, encoding=encoding)
                if not df.empty:
                    return df.to_string()
            except Exception:
                continue
                
        # Last resort fallback
        try:
            with open(file_path, 'rb') as f:
                content = f.read()
                text = content.decode('utf-8', errors='replace')
                return text
        except Exception as e:
            logger.error(f"CSV extraction failed for {file_path.name}: {e}")
            
        return ""
    
    @staticmethod
    def extract_text_from_pptx(file_path: Path) -> str:
        """Extract text from PowerPoint presentations"""
        if pptx is None:
            logger.error("python-pptx library not installed. Cannot extract text from PPTX files.")
            return ""
        
        text = ""
        try:
            prs = pptx.Presentation(file_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
                text += "\n" # Add extra line between slides
        except Exception as e:
            logger.error(f"PPTX extraction error for {file_path.name}: {e}")
        return text
    
    @staticmethod
    def extract_text_from_xlsx(file_path: Path) -> str:
        """Extract text from Excel workbooks"""
        if openpyxl is None:
            logger.error("openpyxl library not installed. Cannot extract text from XLSX files.")
            return ""
        
        text = ""
        try:
            # Try pandas first for a well-formatted output
            try:
                df_dict = pd.read_excel(file_path, sheet_name=None)
                for sheet_name, df in df_dict.items():
                    text += f"Sheet: {sheet_name}\n"
                    text += df.to_string(index=False) + "\n\n"
            except Exception as e:
                logger.warning(f"Pandas Excel extraction failed, falling back to openpyxl: {e}")
                # Fallback to direct openpyxl
                workbook = openpyxl.load_workbook(file_path, data_only=True)
                
                for sheet_name in workbook.sheetnames:
                    sheet = workbook[sheet_name]
                    text += f"Sheet: {sheet_name}\n"
                    
                    for row in sheet.rows:
                        row_text = []
                        for cell in row:
                            value = cell.value
                            if value is not None:
                                row_text.append(str(value))
                            else:
                                row_text.append("")
                        text += "\t".join(row_text) + "\n"
                    text += "\n"  # Add extra line between sheets
        except Exception as e:
            logger.error(f"XLSX extraction error for {file_path.name}: {e}")
        return text
    
    @staticmethod
    def extract_text(file_path: Path) -> str:
        """Extract text from a document based on its extension."""
        ext = file_path.suffix.lower()
            
        if ext == '.pdf':
            return DocumentProcessor.extract_text_from_pdf_with_fallback(file_path)
        elif ext == '.docx' or ext == '.doc':
            return DocumentProcessor.extract_text_from_docx(file_path)
        elif ext == '.txt':
            return DocumentProcessor.extract_text_from_txt(file_path)
        elif ext == '.csv':
            return DocumentProcessor.extract_text_from_csv(file_path)
        elif ext == '.pptx' or ext == '.ppt':
            return DocumentProcessor.extract_text_from_pptx(file_path)
        elif ext == '.xlsx' or ext == '.xls':
            return DocumentProcessor.extract_text_from_xlsx(file_path)
        else:
            logger.warning(f"Unsupported file type: {ext}")
            return ""


class DocumentVectorizer:
    """Process a single document and generate vectors."""

    def __init__(self, config=None):
        self.config = config if config else Config()
        
        # Initialize paths
        self.source_file = Path(self.config.source_file)
        self.vector_db_path = Path(self.config.vector_db_path)
        self.plain_text_path = Path(self.config.plain_text_path)
        self.chunk_path = Path(self.config.chunk_path)
        
        # Make sure output directories exist
        self.vector_db_path.mkdir(parents=True, exist_ok=True)
        self.plain_text_path.mkdir(parents=True, exist_ok=True)
        self.chunk_path.mkdir(parents=True, exist_ok=True)
        
        # Configure chunker
        self.text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=self.config.chunk_size,
            chunk_overlap=self.config.chunk_overlap,
            length_function=len,
        )
        
        # Initialize embeddings model
        self.embeddings = None
        self._load_embedding_model()
        
        # Initialize vector store
        self.vector_store = None
        self._initialize_vector_store()
        
        # Try to load spaCy models
        self.nlp_models = {}
        self._load_language_models()
        
        # Document state
        self.raw_text = None
        self.chunks = []
    
    def _load_embedding_model(self):
        """Load the embedding model."""
        try:
            self.embeddings = HuggingFaceEmbeddings(model_name=self.config.embedding_model)
            logger.info(f"Loaded embedding model: {self.config.embedding_model}")
        except Exception as e:
            logger.error(f"Error loading embedding model: {e}")
            self.embeddings = None
    
    def _load_language_models(self):
        """Load spaCy language models for text processing."""
        if spacy is None:
            logger.warning("spaCy not installed. Language processing functionality limited.")
            return
            
        try:
            # Try to load English model
            if "en" in self.config.languages:
                try:
                    self.nlp_models["en"] = spacy.load("en_core_web_sm")
                    logger.info("Loaded English language model")
                except OSError:
                    logger.warning("English model not found. Run: python -m spacy download en_core_web_sm")
            
            # Try to load Korean model
            if "ko" in self.config.languages:
                try:
                    self.nlp_models["ko"] = spacy.load("ko_core_news_sm")
                    logger.info("Loaded Korean language model")
                except OSError:
                    logger.warning("Korean model not found. Run: python -m spacy download ko_core_news_sm")
            
        except Exception as e:
            logger.error(f"Error loading language models: {e}")
    
    def _initialize_vector_store(self):
        """Initialize or load vector store."""
        vector_store_path = str(self.vector_db_path)
        
        try:
            # Check if vector store already exists
            if (Path(vector_store_path) / "index.faiss").exists():
                logger.info(f"Loading existing vector store from {vector_store_path}")
                self.vector_store = FAISS.load_local(vector_store_path, self.embeddings)
            else:
                logger.info("Vector store does not exist yet")
                self.vector_store = None
        except Exception as e:
            logger.error(f"Error initializing vector store: {e}")
            self.vector_store = None
    
    def _get_file_hash(self, file_path: Path) -> str:
        """Get MD5 hash of file"""
        import hashlib
        
        try:
            hasher = hashlib.md5()
            with open(file_path, 'rb') as f:
                buf = f.read(65536)  # Read in 64kb chunks
                while len(buf) > 0:
                    hasher.update(buf)
                    buf = f.read(65536)
            return hasher.hexdigest()
        except Exception as e:
            logger.warning(f"Cannot hash {file_path.name}: {e}")
            return str(time.time())  # Return timestamp as fallback
    
    def set_source_file(self, file_path: str) -> bool:
        """Set the source file to process."""
        path = Path(file_path)
        if not path.exists():
            logger.error(f"File does not exist: {file_path}")
            return False
        
        self.source_file = path
        self.config.source_file = str(path)
        logger.info(f"Source file set to: {file_path}")
        
        # Reset document state
        self.raw_text = None
        self.chunks = []
        
        return True
    
    def extract_plain_text(self) -> bool:
        """Extract plain text from the source file."""
        if not self.source_file.exists():
            logger.error(f"Source file doesn't exist: {self.source_file}")
            return False
        
        try:
            console.print(f"Extracting text from [bold]{self.source_file.name}[/bold]")
            
            # Extract text based on file type
            self.raw_text = DocumentProcessor.extract_text(self.source_file)
            
            if not self.raw_text or not self.raw_text.strip():
                console.print("[red]No text extracted from the file[/red]")
                return False
            
            # Detect language
            language = self.detect_language(self.raw_text)
            console.print(f"Detected language: [bold]{language}[/bold]")
            
            # Clean text
            self.raw_text = self.clean_text(self.raw_text, language)
            
            # Generate hash for the file
            file_hash = self._get_file_hash(self.source_file)
            
            # Save plain text
            plain_file_name = f"{self.source_file.stem}_{file_hash[:8]}.txt"
            plain_file_path = self.plain_text_path / plain_file_name
            
            with open(plain_file_path, 'w', encoding='utf-8') as f:
                f.write(self.raw_text)
            
            # Update file info
            self.config.current_file_info = {
                'hash': file_hash,
                'language': language,
                'processed_date': datetime.now().isoformat(),
                'chunks': 0,  # Will be updated when chunks are created
                'plain_text_path': str(plain_file_path),
                'chunks_directory': None  # Will be updated when chunks are created
            }
            
            console.print(f"[green]Plain text extracted and saved to:[/green] {plain_file_path}")
            console.print(f"Text length: [bold]{len(self.raw_text)}[/bold] characters")
            
            return True
        
        except Exception as e:
            logger.error(f"Error extracting text: {e}")
            logger.error(traceback.format_exc())
            console.print(f"[red]Error extracting text: {str(e)}[/red]")
            return False
            
    def create_chunks(self, chunk_size: int = None, chunk_overlap: int = None) -> bool:
        """Create text chunks from the extracted plain text."""
        if not self.raw_text:
            console.print("[yellow]No plain text available. Extract text first.[/yellow]")
            return False
        
        # Update chunk settings if provided
        if chunk_size:
            self.config.chunk_size = chunk_size
            self.text_splitter.chunk_size = chunk_size
        
        if chunk_overlap:
            self.config.chunk_overlap = chunk_overlap
            self.text_splitter.chunk_overlap = chunk_overlap
        
        try:
            console.print(f"Creating chunks with size={self.config.chunk_size}, overlap={self.config.chunk_overlap}")
            
            # Create chunks
            self.chunks = self.text_splitter.split_text(self.raw_text)
            
            if not self.chunks:
                console.print("[red]No chunks created from the text[/red]")
                return False
            
            console.print(f"[green]Created {len(self.chunks)} chunks[/green]")
            
            # Save chunks to files
            file_hash = self.config.current_file_info['hash']
            if not file_hash:
                file_hash = self._get_file_hash(self.source_file)
            
            # Create a safer name for directory
            safe_stem = re.sub(r'[^\w\-.]', '_', self.source_file.stem)
            chunks_directory = self.chunk_path / f"{safe_stem}_{file_hash[:8]}"
            chunks_directory.mkdir(exist_ok=True, parents=True)
            
            # Save each chunk to a file
            for i, chunk in enumerate(self.chunks):
                chunk_file_path = chunks_directory / f"chunk_{i:04d}.txt"
                try:
                    with open(chunk_file_path, 'w', encoding='utf-8') as f:
                        f.write(chunk)
                except Exception as e:
                    logger.error(f"Failed to save chunk {i} to file: {e}")
            
            # Update file info
            self.config.current_file_info['chunks'] = len(self.chunks)
            self.config.current_file_info['chunks_directory'] = str(chunks_directory)
            
            console.print(f"Chunks saved to: {chunks_directory}")
            
            return True
        
        except Exception as e:
            logger.error(f"Error creating chunks: {e}")
            logger.error(traceback.format_exc())
            console.print(f"[red]Error creating chunks: {str(e)}[/red]")
            return False
        
    def vectorize_chunks(self, chunk_size: int = None, chunk_overlap: int = None, use_progress: bool = True) -> bool:
        """Create vector embeddings for the chunks."""
        if not self.chunks:
            console.print("[yellow]No chunks available. Create chunks first.[/yellow]")
            return False
        
        if not self.embeddings:
            try:
                self._load_embedding_model()
                if not self.embeddings:
                    console.print("[red]Could not load embedding model[/red]")
                    return False
            except Exception as e:
                console.print(f"[red]Error loading embedding model: {str(e)}[/red]")
                return False
        
        try:
            console.print(f"Vectorizing {len(self.chunks)} chunks with model: [bold]{self.config.embedding_model}[/bold]")
            
            # Create metadata
            metadata = {
                'source': str(self.source_file),
                'filename': self.source_file.name,
                'file_hash': self.config.current_file_info['hash'],
                'language': self.config.current_file_info['language'],
                'processed_date': datetime.now().isoformat(),
                'chunks': len(self.chunks),
                'plain_text_path': self.config.current_file_info['plain_text_path'],
                'chunks_directory': self.config.current_file_info['chunks_directory']
            }
            
            # Create Document objects for each chunk
            documents = []
            for i, chunk in enumerate(self.chunks):
                doc_metadata = metadata.copy()
                doc_metadata['chunk_index'] = i
                chunk_file = Path(self.config.current_file_info['chunks_directory']) / f"chunk_{i:04d}.txt"
                doc_metadata['chunk_file'] = str(chunk_file)
                documents.append(Document(page_content=chunk, metadata=doc_metadata))
            
            # Create or update vector store - with or without progress display
            if use_progress:
                with Progress(
                    SpinnerColumn(),
                    TextColumn("[bold blue]Generating vector embeddings..."),
                    console=console
                ) as progress:
                    progress.add_task("Vectorizing", total=1)
                    
                    if not self.vector_store:
                        self.vector_store = FAISS.from_documents(documents, self.embeddings)
                        console.print(f"[green]Created new vector store with {len(documents)} chunks[/green]")
                    else:
                        self.vector_store.add_documents(documents)
                        console.print(f"[green]Added {len(documents)} chunks to existing vector store[/green]")
            else:
                # No progress display for batch mode
                if not self.vector_store:
                    self.vector_store = FAISS.from_documents(documents, self.embeddings)
                    console.print(f"[green]Created new vector store with {len(documents)} chunks[/green]")
                else:
                    self.vector_store.add_documents(documents)
                    console.print(f"[green]Added {len(documents)} chunks to existing vector store[/green]")
            
            # Save vector store
            self.vector_store.save_local(str(self.vector_db_path))
            
            console.print(f"Vector store saved to: {self.vector_db_path}")
            return True
        
        except Exception as e:
            logger.error(f"Error vectorizing chunks: {e}")
            logger.error(traceback.format_exc())
            console.print(f"[red]Error vectorizing chunks: {str(e)}[/red]")
            return False
       

    def batch_process_folder(self, folder_path: str) -> bool:
        """Process all supported documents in a folder in batch mode."""
        folder = Path(folder_path)
        if not folder.exists() or not folder.is_dir():
            console.print(f"[red]Folder does not exist: {folder_path}[/red]")
            return False
        
        # Scan for supported files
        supported_extensions = ['.pdf', '.docx', '.doc', '.txt', '.csv', '.pptx', '.ppt', '.xlsx', '.xls']
        files_to_process = []
        
        for ext in supported_extensions:
            files_to_process.extend(folder.glob(f"**/*{ext}"))
        
        # Sort files by size (smaller first)
        files_to_process.sort(key=lambda f: f.stat().st_size if f.exists() else 0)
        
        if not files_to_process:
            console.print(f"[yellow]No supported files found in: {folder_path}[/yellow]")
            return False
        
        total_files = len(files_to_process)
        console.print(f"[bold]Found {total_files} files to process[/bold]")
        
        # Ask for batch configuration
        console.print("\n[bold]Batch Processing Settings:[/bold]")
        
        # Ask if custom chunk settings should be used
        use_custom_chunks = Confirm.ask("Use custom chunk settings for all files?", default=False)
        
        chunk_size = self.config.chunk_size
        chunk_overlap = self.config.chunk_overlap
        
        if use_custom_chunks:
            chunk_size = int(Prompt.ask(
                "Chunk size (characters)",
                default=str(self.config.chunk_size)
            ))
            
            chunk_overlap = int(Prompt.ask(
                "Chunk overlap (characters)",
                default=str(self.config.chunk_overlap)
            ))
        
        # Create progress tracking
        completed_files = 0
        success_count = 0
        failed_count = 0
        skipped_count = 0
        
        # Create statistics for reporting
        stats = {
            'by_extension': {},
            'extract_failures': 0,
            'chunk_failures': 0,
            'vector_failures': 0,
            'total_chunks': 0,
            'total_characters': 0
        }
        
        # Initialize extension stats
        for ext in supported_extensions:
            stats['by_extension'][ext] = {'count': 0, 'success': 0, 'failed': 0, 'skipped': 0}
        
        # Create a batch processing log
        batch_log_path = folder / "batch_processing_log.json"
        processed_files = {}
        
        # Load existing log if available for resuming
        if batch_log_path.exists():
            try:
                with open(batch_log_path, 'r', encoding='utf-8') as f:
                    processed_files = json.load(f)
                console.print(f"[green]Found existing batch log with {len(processed_files)} processed files[/green]")
                resume = Confirm.ask("Resume previous batch processing?", default=True)
                if not resume:
                    processed_files = {}
            except Exception as e:
                console.print(f"[yellow]Could not load batch log: {str(e)}. Starting fresh.[/yellow]")
                processed_files = {}
        
        start_time = time.time()
        
        # Process each file
        with Progress(
            SpinnerColumn(),
            TextColumn("[bold blue]{task.description}"),
            BarColumn(),
            TextColumn("[bold]{task.completed}/{task.total}"),
            TextColumn("[green]{task.fields[status]}"),
            console=console
        ) as progress:
            batch_task = progress.add_task(
                "[BATCH PROCESSING]", 
                total=total_files,
                status="Starting...",
                start=True
            )
            
            for file_idx, file_path in enumerate(files_to_process, 1):
                file_str = str(file_path)
                file_ext = file_path.suffix.lower()
                
                # Update extension stats
                if file_ext in stats['by_extension']:
                    stats['by_extension'][file_ext]['count'] += 1
                
                file_status = "Processing"
                progress.update(batch_task, description=f"[BATCH PROCESSING] {file_path.name}", status=file_status)
                
                # Check if file was already processed
                file_hash = self._get_file_hash(file_path)
                if file_str in processed_files and processed_files[file_str].get('hash') == file_hash:
                    console.print(f"[yellow]Skipping previously processed file: {file_path.name}[/yellow]")
                    skipped_count += 1
                    
                    # Update extension stats
                    if file_ext in stats['by_extension']:
                        stats['by_extension'][file_ext]['skipped'] += 1
                    
                    completed_files += 1
                    progress.update(batch_task, completed=completed_files, 
                                status=f"Skipped: {completed_files}/{total_files} files | Success: {success_count} | Failed: {failed_count} | Skipped: {skipped_count}")
                    continue
                
                # Process the file
                try:
                    # Set source file
                    file_status = "Extracting text"
                    progress.update(batch_task, status=file_status)
                    
                    # Use a temporary copy of the file info to avoid overwriting
                    orig_file_info = self.config.current_file_info.copy()
                    self.set_source_file(file_str)
                    
                    # Extract text
                    if not self.extract_plain_text():
                        console.print(f"[red]Failed to extract text from: {file_path.name}[/red]")
                        failed_count += 1
                        stats['extract_failures'] += 1
                        
                        # Update extension stats
                        if file_ext in stats['by_extension']:
                            stats['by_extension'][file_ext]['failed'] += 1
                        
                        completed_files += 1
                        progress.update(batch_task, completed=completed_files, 
                                    status=f"Failed: {completed_files}/{total_files} files | Success: {success_count} | Failed: {failed_count} | Skipped: {skipped_count}")
                        
                        # Restore original file info
                        self.config.current_file_info = orig_file_info
                        continue
                    
                    # Update stats with text length
                    stats['total_characters'] += len(self.raw_text or "")
                    
                    # Create chunks
                    file_status = "Creating chunks"
                    progress.update(batch_task, status=file_status)
                    
                    if not self.create_chunks(chunk_size, chunk_overlap):
                        console.print(f"[red]Failed to create chunks for: {file_path.name}[/red]")
                        failed_count += 1
                        stats['chunk_failures'] += 1
                        
                        # Update extension stats
                        if file_ext in stats['by_extension']:
                            stats['by_extension'][file_ext]['failed'] += 1
                        
                        completed_files += 1
                        progress.update(batch_task, completed=completed_files, 
                                    status=f"Failed: {completed_files}/{total_files} files | Success: {success_count} | Failed: {failed_count} | Skipped: {skipped_count}")
                        
                        # Restore original file info
                        self.config.current_file_info = orig_file_info
                        continue
                    
                    # Update stats with chunk count
                    stats['total_chunks'] += len(self.chunks)
                    
                    # Vectorize chunks - use False for progress to avoid nesting progress bars
                    file_status = "Vectorizing"
                    progress.update(batch_task, status=file_status)
                    
                    if not self.vectorize_chunks(use_progress=False):
                        console.print(f"[red]Failed to vectorize chunks for: {file_path.name}[/red]")
                        failed_count += 1
                        stats['vector_failures'] += 1
                        
                        # Update extension stats
                        if file_ext in stats['by_extension']:
                            stats['by_extension'][file_ext]['failed'] += 1
                        
                        completed_files += 1
                        progress.update(batch_task, completed=completed_files, 
                                    status=f"Failed: {completed_files}/{total_files} files | Success: {success_count} | Failed: {failed_count} | Skipped: {skipped_count}")
                        
                        # Restore original file info
                        self.config.current_file_info = orig_file_info
                        continue
                    
                    # Record successful processing
                    processed_files[file_str] = {
                        'hash': file_hash,
                        'language': self.config.current_file_info.get('language', 'unknown'),
                        'processed_date': datetime.now().isoformat(),
                        'chunks': self.config.current_file_info.get('chunks', 0),
                        'plain_text_path': self.config.current_file_info.get('plain_text_path'),
                        'chunks_directory': self.config.current_file_info.get('chunks_directory')
                    }
                    
                    # Save batch log periodically
                    if file_idx % 5 == 0 or file_idx == total_files:
                        try:
                            with open(batch_log_path, 'w', encoding='utf-8') as f:
                                json.dump(processed_files, f, indent=2, ensure_ascii=False)
                        except Exception as e:
                            console.print(f"[yellow]Could not save batch log: {str(e)}[/yellow]")
                    
                    success_count += 1
                    
                    # Update extension stats
                    if file_ext in stats['by_extension']:
                        stats['by_extension'][file_ext]['success'] += 1
                    
                    file_status = "Success"
                    
                except Exception as e:
                    console.print(f"[red]Error processing {file_path.name}: {str(e)}[/red]")
                    logger.error(f"Error in batch processing file {file_path}: {e}")
                    logger.error(traceback.format_exc())
                    failed_count += 1
                    
                    # Update extension stats
                    if file_ext in stats['by_extension']:
                        stats['by_extension'][file_ext]['failed'] += 1
                    
                    file_status = "Error"
                
                # Update progress
                completed_files += 1
                progress.update(
                    batch_task, 
                    completed=completed_files,
                    status=f"{file_status}: {completed_files}/{total_files} files | Success: {success_count} | Failed: {failed_count} | Skipped: {skipped_count}"
                )
                
                # Restore original file info after processing
                self.config.current_file_info = orig_file_info
        
        # Calculate processing time
        elapsed_time = time.time() - start_time
        hours, remainder = divmod(elapsed_time, 3600)
        minutes, seconds = divmod(remainder, 60)
        time_str = f"{int(hours)}h {int(minutes)}m {int(seconds)}s"
        
        # Show final statistics
        console.print("\n[bold blue]BATCH PROCESSING COMPLETED[/bold blue]")
        console.print(f"[bold]Total files processed:[/bold] {total_files}")
        console.print(f"[bold]Successful:[/bold] {success_count}")
        console.print(f"[bold]Failed:[/bold] {failed_count}")
        console.print(f"[bold]Skipped:[/bold] {skipped_count}")
        console.print(f"[bold]Total processing time:[/bold] {time_str}")
        
        # Show detailed statistics
        console.print("\n[bold]Processing Statistics:[/bold]")
        console.print(f"Total text extracted: {stats['total_characters']:,} characters")
        console.print(f"Total chunks created: {stats['total_chunks']:,}")
        console.print(f"Extract failures: {stats['extract_failures']}")
        console.print(f"Chunk failures: {stats['chunk_failures']}")
        console.print(f"Vector failures: {stats['vector_failures']}")
        
        # Show extension statistics in a table
        console.print("\n[bold]File Type Statistics:[/bold]")
        table = Table(title="Processing by File Type")
        table.add_column("Extension", style="cyan")
        table.add_column("Total", style="white", justify="right")
        table.add_column("Success", style="green", justify="right")
        table.add_column("Failed", style="red", justify="right")
        table.add_column("Skipped", style="yellow", justify="right")
        
        for ext, data in sorted(stats['by_extension'].items()):
            if data['count'] > 0:  # Only show extensions that had files
                table.add_row(
                    ext,
                    str(data['count']),
                    str(data['success']),
                    str(data['failed']),
                    str(data['skipped'])
                )
        
        console.print(table)
        
        # Save final batch log
        try:
            with open(batch_log_path, 'w', encoding='utf-8') as f:
                json.dump(processed_files, f, indent=2, ensure_ascii=False)
            console.print(f"[green]Batch processing log saved to: {batch_log_path}[/green]")
        except Exception as e:
            console.print(f"[yellow]Could not save batch log: {str(e)}[/yellow]")
        
        return True


    def show_plain_text(self):
        """Display the extracted plain text."""
        if not self.raw_text:
            # Try to load from file if available
            if self.config.current_file_info.get('plain_text_path'):
                try:
                    plain_path = Path(self.config.current_file_info['plain_text_path'])
                    if plain_path.exists():
                        with open(plain_path, 'r', encoding='utf-8') as f:
                            self.raw_text = f.read()
                    else:
                        console.print("[yellow]Plain text file not found[/yellow]")
                        return
                except Exception as e:
                    console.print(f"[red]Error loading plain text: {str(e)}[/red]")
                    return
            else:
                console.print("[yellow]No plain text available. Extract text first.[/yellow]")
                return
        
        # Display text statistics
        text_length = len(self.raw_text)
        line_count = self.raw_text.count('\n') + 1
        word_count = len(self.raw_text.split())
        
        console.print(f"[bold]Plain Text from:[/bold] {self.source_file.name}")
        console.print(f"Length: {text_length} characters")
        console.print(f"Lines: {line_count}")
        console.print(f"Words: {word_count}")
        console.print(f"Language: {self.config.current_file_info.get('language', 'unknown')}")
        
        # Display text preview
        preview_length = min(1000, len(self.raw_text))
        preview_text = self.raw_text[:preview_length]
        if len(self.raw_text) > preview_length:
            preview_text += "\n...(truncated)..."
        
        console.print("\n[bold]Text Preview:[/bold]")
        console.print(Panel(preview_text, title="Plain Text Preview", width=100))
        
        # Ask if user wants to view the full text
        if len(self.raw_text) > preview_length and Confirm.ask("View full text?", default=False):
            # Use syntax highlighting for plain text
            syntax = Syntax(self.raw_text, "text", theme="monokai", line_numbers=True, word_wrap=True)
            console.print(syntax)
    
    def show_chunks(self):
        """Display the created chunks."""
        if not self.chunks:
            # Try to load chunks from files if available
            chunks_dir = self.config.current_file_info.get('chunks_directory')
            if chunks_dir and Path(chunks_dir).exists():
                try:
                    chunk_files = sorted(Path(chunks_dir).glob("chunk_*.txt"))
                    self.chunks = []
                    for chunk_file in chunk_files:
                        with open(chunk_file, 'r', encoding='utf-8') as f:
                            self.chunks.append(f.read())
                    console.print(f"[green]Loaded {len(self.chunks)} chunks from {chunks_dir}[/green]")
                except Exception as e:
                    console.print(f"[red]Error loading chunks: {str(e)}[/red]")
                    return
            else:
                console.print("[yellow]No chunks available. Create chunks first.[/yellow]")
                return
        
        # Display chunks table
        console.print(f"[bold]Chunks from:[/bold] {self.source_file.name}")
        console.print(f"Total Chunks: {len(self.chunks)}")
        console.print(f"Chunk Size: {self.config.chunk_size}")
        console.print(f"Chunk Overlap: {self.config.chunk_overlap}")
        
        # Create a table with all chunks
        table = Table(title="Document Chunks")
        table.add_column("#", style="cyan", justify="right")
        table.add_column("Length", style="green", justify="right")
        table.add_column("Preview", style="white")
        
        for i, chunk in enumerate(self.chunks):
            # Get the first line or two for preview
            preview_lines = chunk.split('\n')[:2]
            preview = '\n'.join(preview_lines)
            if len(preview) > 80:
                preview = preview[:77] + "..."
            
            table.add_row(
                str(i+1), 
                str(len(chunk)),
                preview
            )
        
        console.print(table)
        
        # Ask if user wants to view a specific chunk
        if Confirm.ask("View a specific chunk?", default=False):
            chunk_num = int(Prompt.ask(
                f"Enter chunk number (1-{len(self.chunks)})",
                default="1"
            ))
            
            if 1 <= chunk_num <= len(self.chunks):
                chunk_text = self.chunks[chunk_num-1]
                console.print(Panel(
                    chunk_text,
                    title=f"Chunk {chunk_num}/{len(self.chunks)}",
                    width=100
                ))
    
    def query_vector_store(self):
        """Query the vector store for similar chunks."""
        if not self.vector_store:
            self._initialize_vector_store()
            
        if not self.vector_store:
            console.print("[yellow]No vector store available. Vectorize chunks first.[/yellow]")
            return
        
        # Get query from user
        query = Prompt.ask("Enter your query")
        
        if not query.strip():
            console.print("[yellow]Empty query[/yellow]")
            return
        
        # Get number of results to show
        k = int(Prompt.ask("Number of results to show", default="3"))
        
        try:
            console.print(f"\n[bold]Searching for:[/bold] {query}")
            
            # Perform similarity search
            with Progress(
                SpinnerColumn(),
                TextColumn("[bold blue]Searching vector store..."),
                console=console
            ) as progress:
                task = progress.add_task("Searching", total=1)
                results = self.vector_store.similarity_search_with_score(query, k=k)
                progress.update(task, completed=1)
            
            if not results:
                console.print("[yellow]No results found[/yellow]")
                return
            
            # Display results
            result_table = Table(title=f"Top {len(results)} Results")
            result_table.add_column("Score", style="cyan")
            result_table.add_column("Source", style="green")
            result_table.add_column("Chunk", style="magenta")
            result_table.add_column("Content", style="white")
            
            for i, (doc, score) in enumerate(results, 1):
                # Format content text for display (truncate if too long)
                content = doc.page_content.strip()
                if len(content) > 200:
                    content = content[:197] + "..."
                
                # Format score (convert distance to similarity score)
                formatted_score = f"{1.0 - score:.4f}"
                
                # Get metadata
                source = Path(doc.metadata.get('source', 'Unknown')).name
                chunk_index = doc.metadata.get('chunk_index', 'Unknown')
                
                result_table.add_row(
                    formatted_score,
                    source,
                    str(chunk_index),
                    content
                )
            
            console.print(result_table)
            
            # Ask if user wants to see detailed content
            if Confirm.ask("Show full content of a specific result?", default=False):
                result_num = int(Prompt.ask(
                    f"Enter result number (1-{len(results)})",
                    default="1"
                ))
                
                if 1 <= result_num <= len(results):
                    doc, score = results[result_num - 1]
                    
                    # Display full content in a panel
                    metadata = doc.metadata
                    metadata_str = "\n".join([f"{k}: {v}" for k, v in metadata.items() 
                                             if k not in ['chunk_file', 'plain_text_path', 'chunks_directory']])
                    
                    console.print(Panel(
                        f"[bold]Metadata:[/bold]\n{metadata_str}\n\n[bold]Content:[/bold]\n{doc.page_content}",
                        title=f"Result {result_num} Details",
                        expand=False
                        ))
                else:
                    console.print("[yellow]Invalid result number[/yellow]")
        
        except Exception as e:
            logger.error(f"Error querying vector store: {e}")
            logger.error(traceback.format_exc())
            console.print(f"[red]Error querying vector store: {str(e)}[/red]")
    
    def detect_language(self, text: str) -> str:
        """Detect the language of the text."""
        if not text or len(text.strip()) < 10:
            return "en"  # Default to English for empty or very short texts
        
        # Simple heuristic: count Korean and English characters
        ko_char_count = len(re.findall(r'[가-힣]', text))
        en_char_count = len(re.findall(r'[a-zA-Z]', text))
        
        if ko_char_count > en_char_count:
            return "ko"
        else:
            return "en"
    
    def clean_text(self, text: str, language: str = "en") -> str:
        """Clean and normalize text based on detected language."""
        if not text:
            return ""
        
        # Basic cleaning
        text = re.sub(r'\s+', ' ', text)  # Normalize whitespace
        text = text.strip()
        
        # Use spaCy for advanced processing if model is available
        if spacy is not None and language in self.nlp_models:
            nlp = self.nlp_models[language]
            # Process with spaCy, but limit document size to avoid memory issues
            max_length = 100000  # 100k characters max
            if len(text) > max_length:
                logger.warning(f"Text too long for spaCy processing, truncating to {max_length} characters")
                text = text[:max_length]
            
            try:
                doc = nlp(text)
                # Basic sentence normalization
                processed_text = " ".join([sent.text.strip() for sent in doc.sents])
                return processed_text
            except Exception as e:
                logger.error(f"Error in spaCy processing: {e}")
                return text
        
        return text
    
    def clear_vector_database(self) -> bool:
        """Clear the vector database."""
        try:
            vector_path = self.vector_db_path
            
            # Find all vector store related files
            vector_files = list(vector_path.glob("*.faiss")) + list(vector_path.glob("*.pkl"))
            
            if not vector_files:
                console.print("[yellow]No vector database files found[/yellow]")
                return False
            
            # Delete each file
            for file in vector_files:
                file.unlink()
                console.print(f"Deleted: {file}")
            
            # Reset vector store
            self.vector_store = None
            console.print("[green]Vector database successfully cleared[/green]")
            return True
        
        except Exception as e:
            logger.error(f"Error clearing vector database: {e}")
            console.print(f"[red]Error clearing vector database: {str(e)}[/red]")
            return False
    
    def clear_chunks(self) -> bool:
        """Clear the chunk files."""
        try:
            # Check if there's a specific chunk directory to delete
            if self.config.current_file_info.get('chunks_directory'):
                chunks_dir = Path(self.config.current_file_info['chunks_directory'])
                if chunks_dir.exists():
                    # Delete all files in the directory
                    for chunk_file in chunks_dir.glob("chunk_*.txt"):
                        chunk_file.unlink()
                        
                    # Try to remove the directory
                    chunks_dir.rmdir()
                    console.print(f"[green]Cleared chunks directory: {chunks_dir}[/green]")
                    
                    # Update file info
                    self.config.current_file_info['chunks_directory'] = None
                    self.config.current_file_info['chunks'] = 0
                    
                    # Reset chunks
                    self.chunks = []
                    return True
                else:
                    console.print("[yellow]Chunks directory not found[/yellow]")
            
            # If no specific directory or couldn't delete it, ask if all chunks should be deleted
            if Confirm.ask("Delete all chunk directories?", default=False):
                # Delete all contents of the chunks path
                chunk_path = self.chunk_path
                deleted_count = 0
                
                for item in chunk_path.glob("*"):
                    if item.is_dir():
                        for chunk_file in item.glob("*"):
                            chunk_file.unlink()
                        item.rmdir()
                        deleted_count += 1
                
                if deleted_count > 0:
                    console.print(f"[green]Deleted {deleted_count} chunk directories[/green]")
                    # Reset chunks
                    self.chunks = []
                    self.config.current_file_info['chunks_directory'] = None
                    self.config.current_file_info['chunks'] = 0
                    return True
                else:
                    console.print("[yellow]No chunk directories found[/yellow]")
                    return False
            
            return False
        
        except Exception as e:
            logger.error(f"Error clearing chunks: {e}")
            console.print(f"[red]Error clearing chunks: {str(e)}[/red]")
            return False
    
    def clear_plain_texts(self) -> bool:
        """Clear the plain text files."""
        try:
            # Check if there's a specific plain text file to delete
            if self.config.current_file_info.get('plain_text_path'):
                plain_file = Path(self.config.current_file_info['plain_text_path'])
                if plain_file.exists():
                    plain_file.unlink()
                    console.print(f"[green]Deleted plain text file: {plain_file}[/green]")
                    
                    # Update file info
                    self.config.current_file_info['plain_text_path'] = None
                    
                    # Reset raw text
                    self.raw_text = None
                    return True
                else:
                    console.print("[yellow]Plain text file not found[/yellow]")
            
            # If no specific file or couldn't delete it, ask if all plain texts should be deleted
            if Confirm.ask("Delete all plain text files?", default=False):
                # Delete all plain text files
                plain_path = self.plain_text_path
                deleted_count = 0
                
                for plain_file in plain_path.glob("*.txt"):
                    plain_file.unlink()
                    deleted_count += 1
                
                if deleted_count > 0:
                    console.print(f"[green]Deleted {deleted_count} plain text files[/green]")
                    # Reset raw text
                    self.raw_text = None
                    self.config.current_file_info['plain_text_path'] = None
                    return True
                else:
                    console.print("[yellow]No plain text files found[/yellow]")
                    return False
            
            return False
        
        except Exception as e:
            logger.error(f"Error clearing plain texts: {e}")
            console.print(f"[red]Error clearing plain texts: {str(e)}[/red]")
            return False
    
    def configure_settings(self):
        """Configure the vectorizer settings."""
        console.print("[bold blue]Configuration Settings[/bold blue]")
        
        # Chunking settings
        console.print("\n[bold]Text Chunking Settings:[/bold]")
        
        self.config.chunk_size = int(Prompt.ask(
            "Chunk size (characters)",
            default=str(self.config.chunk_size)
        ))
        
        self.config.chunk_overlap = int(Prompt.ask(
            "Chunk overlap (characters)",
            default=str(self.config.chunk_overlap)
        ))
        
        # Update the text splitter with new settings
        self.text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=self.config.chunk_size,
            chunk_overlap=self.config.chunk_overlap,
            length_function=len,
        )
        
        # Language settings
        console.print("\n[bold]Language Settings:[/bold]")
        console.print("Select languages to support:")
        console.print("1. English only")
        console.print("2. Korean only")
        console.print("3. Both English and Korean")
        
        lang_choice = Prompt.ask(
            "Choose language option",
            choices=["1", "2", "3"],
            default="3"
        )
        
        if lang_choice == "1":
            self.config.languages = ["en"]
        elif lang_choice == "2":
            self.config.languages = ["ko"]
        else:
            self.config.languages = ["en", "ko"]
        
        # Reload language models if needed
        self._load_language_models()
        
        # Embedding model
        console.print("\n[bold]Embedding Model:[/bold]")
        console.print("1. all-MiniLM-L6-v2 (fast, lower quality)")
        console.print("2. msmarco-distilbert-base-v4 (balanced)")
        console.print("3. all-mpnet-base-v2 (slower, high quality)")
        
        model_choice = Prompt.ask(
            "Choose embedding model",
            choices=["1", "2", "3"],
            default="1"
        )
        
        old_model = self.config.embedding_model
        if model_choice == "1":
            self.config.embedding_model = "all-MiniLM-L6-v2"
        elif model_choice == "2":
            self.config.embedding_model = "msmarco-distilbert-base-v4"
        elif model_choice == "3":
            self.config.embedding_model = "all-mpnet-base-v2"
        
        # Check if model changed
        if old_model != self.config.embedding_model:
            console.print(f"[yellow]Embedding model changed from {old_model} to {self.config.embedding_model}[/yellow]")
            console.print("[yellow]Note: Existing vector store will be incompatible with the new model[/yellow]")
            
            # Reload embedding model
            self._load_embedding_model()
            
            # Reset vector store if model changed
            self.vector_store = None
        
        # Directory settings
        console.print("\n[bold]Directory Settings:[/bold]")
        
        new_plain_path = Prompt.ask(
            "Plain text directory path",
            default=str(self.plain_text_path)
        )
        
        new_chunk_path = Prompt.ask(
            "Chunk directory path",
            default=str(self.chunk_path)
        )
        
        new_vector_path = Prompt.ask(
            "Vector database path",
            default=str(self.vector_db_path)
        )
        
        # Update paths if changed
        if new_plain_path != str(self.plain_text_path):
            self.config.plain_text_path = new_plain_path
            self.plain_text_path = Path(new_plain_path)
            self.plain_text_path.mkdir(parents=True, exist_ok=True)
        
        if new_chunk_path != str(self.chunk_path):
            self.config.chunk_path = new_chunk_path
            self.chunk_path = Path(new_chunk_path)
            self.chunk_path.mkdir(parents=True, exist_ok=True)
        
        if new_vector_path != str(self.vector_db_path):
            self.config.vector_db_path = new_vector_path
            self.vector_db_path = Path(new_vector_path)
            self.vector_db_path.mkdir(parents=True, exist_ok=True)
            
            # Reset vector store if path changed
            self.vector_store = None
            self._initialize_vector_store()
        
        # Save configuration
        if self.config.save():
            console.print("[green]Configuration saved successfully![/green]")
        else:
            console.print("[red]Failed to save configuration[/red]")


class DocVectorizerApp:
    """Main application for document vectorization."""
    
    def __init__(self):
        self.config = Config()
        
        # Try to load configuration
        try:
            self.config.load()
        except Exception as e:
            logger.warning(f"Could not load configuration: {e}")
            # If loading fails, use default configuration
            pass
        
        # Initialize vectorizer
        self.vectorizer = DocumentVectorizer(self.config)
        
        # Check for required directories
        self._check_directories()
    
    def _check_directories(self):
        """Ensure required directories exist."""
        os.makedirs(self.config.plain_text_path, exist_ok=True)
        os.makedirs(self.config.chunk_path, exist_ok=True)
        os.makedirs(self.config.vector_db_path, exist_ok=True)
        
    def _select_source_file(self):
        """Dialog to select a source file."""
        console.print("[bold blue]Select Source File[/bold blue]")
        
        current_file = self.config.source_file
        console.print(f"Current file: [bold]{current_file}[/bold]")
        
        options = [
            "Enter file path manually",
            "Scan current directory for files",
            "Scan a directory for files with extension"
        ]
        
        console.print("\nOptions:")
        for i, opt in enumerate(options, 1):
            console.print(f"{i}. {opt}")
        
        choice = Prompt.ask("Select an option", choices=["1", "2", "3"], default="1")
        
        if choice == "1":
            # Manual path entry
            file_path = Prompt.ask("Enter file path", default=current_file)
            return self.vectorizer.set_source_file(file_path)
        
        elif choice == "2":
            # Scan current directory
            console.print("\nScanning current directory for supported files...")
            
            supported_extensions = ['.pdf', '.docx', '.doc', '.txt', '.csv', '.pptx', '.ppt', '.xlsx', '.xls']
            files = []
            
            for ext in supported_extensions:
                files.extend(Path.cwd().glob(f"*{ext}"))
            
            if not files:
                console.print("[yellow]No supported files found in current directory[/yellow]")
                return False
            
            # Sort files by name
            files.sort(key=lambda f: f.name)
            
            # Display file list
            console.print("\nAvailable files:")
            for i, file in enumerate(files, 1):
                console.print(f"{i}. {file.name} [{file.suffix.lower()}]")
            
            # Ask for selection
            file_choice = Prompt.ask(
                "Select a file by number",
                choices=[str(i) for i in range(1, len(files) + 1)],
                default="1"
            )
            
            selected_file = files[int(file_choice) - 1]
            console.print(f"Selected: [bold]{selected_file}[/bold]")
            
            return self.vectorizer.set_source_file(str(selected_file))
            
        elif choice == "3":
            # Scan directory for files with specific extension
            scan_dir = Prompt.ask("Enter directory path to scan", default=str(Path.cwd()))
            
            # Default extension is .log as requested
            extension = Prompt.ask("Enter file extension (include the dot, e.g. .log)", default=".log")
            
            try:
                scan_path = Path(scan_dir)
                if not scan_path.exists() or not scan_path.is_dir():
                    console.print(f"[red]Invalid directory: {scan_dir}[/red]")
                    return False
                
                # Scan for files with the specified extension
                console.print(f"\nScanning {scan_dir} for *{extension} files...")
                files = list(scan_path.glob(f"*{extension}"))
                
                if not files:
                    console.print(f"[yellow]No {extension} files found in {scan_dir}[/yellow]")
                    return False
                
                # Sort files by name
                files.sort(key=lambda f: f.name)
                
                # Display file list
                console.print("\nAvailable files:")
                for i, file in enumerate(files, 1):
                    console.print(f"{i}. {file.name}")
                
                # Ask for selection
                file_choice = Prompt.ask(
                    "Select a file by number",
                    choices=[str(i) for i in range(1, len(files) + 1)],
                    default="1"
                )
                
                selected_file = files[int(file_choice) - 1]
                console.print(f"Selected: [bold]{selected_file}[/bold]")
                
                return self.vectorizer.set_source_file(str(selected_file))
                
            except Exception as e:
                console.print(f"[red]Error scanning directory: {str(e)}[/red]")
                logger.error(f"Error scanning directory: {e}")
                return False


    def run(self):
        """Run the main application loop."""
        try:
            while True:
                console.clear()
                console.print("[bold blue]Document Vectorizer[/bold blue]")
                console.print("[yellow]=========================[/yellow]")
                console.print("A step-by-step tool for document vectorization")
                
                # Show current file
                console.print(f"\n[bold]Current Source File:[/bold] {self.vectorizer.source_file.name}")
                
                # Check document processing states
                has_plain_text = self.vectorizer.raw_text is not None
                if not has_plain_text and self.config.current_file_info.get('plain_text_path'):
                    plain_path = Path(self.config.current_file_info['plain_text_path'])
                    has_plain_text = plain_path.exists()
                
                has_chunks = len(self.vectorizer.chunks) > 0
                if not has_chunks and self.config.current_file_info.get('chunks_directory'):
                    chunks_dir = Path(self.config.current_file_info['chunks_directory'])
                    has_chunks = chunks_dir.exists() and len(list(chunks_dir.glob("chunk_*.txt"))) > 0
                
                has_vectors = self.vectorizer.vector_store is not None
                
                # Status indicators
                console.print(f"[bold]Processing Status:[/bold]")
                console.print(f"Plain Text: {'[green]✓[/green]' if has_plain_text else '[red]✗[/red]'}")
                console.print(f"Chunks: {'[green]✓[/green]' if has_chunks else '[red]✗[/red]'}")
                console.print(f"Vectors: {'[green]✓[/green]' if has_vectors else '[red]✗[/red]'}")
                
                # Show menu
                console.print("\n[bold]Menu:[/bold]")
                console.print("1. Set source file (PDF, DOC, DOCX, TXT, CSV, PPT, PPTX, XLS, XLSX)")
                console.print("2. Extract plain text from source file")
                console.print("3. Create chunks from plain text")
                console.print("4. Vectorize chunks")
                console.print("5. [bold yellow]BATCH PROCESS FOLDER[/bold yellow]")
                console.print("6. Show plain text")
                console.print("7. Show chunks")
                console.print("8. Vector query")
                console.print("9. Clear vector database")
                console.print("10. Clear chunks")
                console.print("11. Clear plain texts")
                console.print("12. Configuration")
                console.print("13. Exit")
                
                choice = Prompt.ask("Select an option", choices=[str(i) for i in range(1, 14)], default="1")
                
                console.clear()
                
                if choice == "1":
                    # Set source file
                    self._select_source_file()
                    Prompt.ask("Press Enter to continue")
                
                elif choice == "2":
                    # Extract plain text
                    console.print("[bold blue]Extract Plain Text[/bold blue]")
                    self.vectorizer.extract_plain_text()
                    Prompt.ask("Press Enter to continue")
                
                elif choice == "3":
                    # Create chunks
                    console.print("[bold blue]Create Chunks[/bold blue]")
                    
                    # Ask for custom chunk settings
                    use_custom = Confirm.ask("Use custom chunk settings?", default=False)
                    
                    chunk_size = None
                    chunk_overlap = None
                    
                    if use_custom:
                        chunk_size = int(Prompt.ask(
                            "Chunk size (characters)",
                            default=str(self.config.chunk_size)
                        ))
                        
                        chunk_overlap = int(Prompt.ask(
                            "Chunk overlap (characters)",
                            default=str(self.config.chunk_overlap)
                        ))
                    
                    self.vectorizer.create_chunks(chunk_size, chunk_overlap)
                    Prompt.ask("Press Enter to continue")
                
                elif choice == "4":
                    # Vectorize chunks
                    console.print("[bold blue]Vectorize Chunks[/bold blue]")
                    self.vectorizer.vectorize_chunks()
                    Prompt.ask("Press Enter to continue")
                
                elif choice == "5":
                    # Batch process folder
                    console.print("[bold blue]BATCH PROCESS FOLDER[/bold blue]")
                    console.print("[yellow]This will extract text, create chunks, and vectorize all supported documents in a folder.[/yellow]")
                    
                    folder_path = Prompt.ask("Enter folder path to process", default=str(Path.cwd()))
                    
                    if Confirm.ask(f"Process all documents in: {folder_path}?", default=True):
                        self.vectorizer.batch_process_folder(folder_path)
                    
                    Prompt.ask("Press Enter to continue")
                
                elif choice == "6":
                    # Show plain text
                    console.print("[bold blue]Plain Text View[/bold blue]")
                    self.vectorizer.show_plain_text()
                    Prompt.ask("Press Enter to continue")
                
                elif choice == "7":
                    # Show chunks
                    console.print("[bold blue]Chunks View[/bold blue]")
                    self.vectorizer.show_chunks()
                    Prompt.ask("Press Enter to continue")
                
                elif choice == "8":
                    # Vector query
                    console.print("[bold blue]Vector Query[/bold blue]")
                    self.vectorizer.query_vector_store()
                    Prompt.ask("Press Enter to continue")
                
                elif choice == "9":
                    # Clear vector database
                    console.print("[bold blue]Clear Vector Database[/bold blue]")
                    if Confirm.ask("Are you sure you want to clear the vector database?", default=False):
                        self.vectorizer.clear_vector_database()
                    Prompt.ask("Press Enter to continue")
                
                elif choice == "10":
                    # Clear chunks
                    console.print("[bold blue]Clear Chunks[/bold blue]")
                    if Confirm.ask("Are you sure you want to clear the chunks?", default=False):
                        self.vectorizer.clear_chunks()
                    Prompt.ask("Press Enter to continue")
                
                elif choice == "11":
                    # Clear plain texts
                    console.print("[bold blue]Clear Plain Texts[/bold blue]")
                    if Confirm.ask("Are you sure you want to clear the plain texts?", default=False):
                        self.vectorizer.clear_plain_texts()
                    Prompt.ask("Press Enter to continue")
                
                elif choice == "12":
                    # Configuration
                    console.print("[bold blue]Configuration[/bold blue]")
                    self.vectorizer.configure_settings()
                    Prompt.ask("Press Enter to continue")
                
                elif choice == "13":
                    # Exit
                    break
        
        except KeyboardInterrupt:
            console.print("\n[yellow]Application terminated by user[/yellow]")
        except Exception as e:
            console.print(f"[red]An unexpected error occurred: {str(e)}[/red]")
            logger.error(f"Unexpected error: {str(e)}\n{traceback.format_exc()}")
        finally:
            console.print("[green]Exiting Document Vectorizer. Goodbye![/green]")
            


def show_welcome_message():
    """Display welcome message and instructions."""
    console.print("[bold blue]Welcome to Document Vectorizer![/bold blue]")
    console.print("[yellow]==============================[/yellow]")
    console.print("""
This tool helps you process a single document through the steps of:
1. Extracting plain text
2. Creating chunks of text
3. Generating vector embeddings

Use the menu to navigate through each step of the process.
    """)
    Prompt.ask("Press Enter to start")


def check_dependencies():
    """Check if required dependencies are installed."""
    missing_deps = []
    
    # Check for document processing libraries
    try:
        import docx
    except ImportError:
        missing_deps.append("python-docx")
    
    try:
        import PyPDF2
    except ImportError:
        missing_deps.append("PyPDF2")
    
    # Check for ML/Vector libraries
    try:
        import langchain_huggingface
    except ImportError:
        missing_deps.append("langchain-huggingface")
    
    try:
        from langchain_community.vectorstores import FAISS
    except ImportError:
        missing_deps.append("langchain-community")
    
    try:
        import langchain
    except ImportError:
        missing_deps.append("langchain")
    
    # Check for optional dependencies
    optional_missing = []
    
    try:
        import pptx
    except ImportError:
        optional_missing.append("python-pptx")
    
    try:
        import openpyxl
    except ImportError:
        optional_missing.append("openpyxl")
    
    try:
        import spacy
    except ImportError:
        optional_missing.append("spacy")


    # Check for OCR dependencies
    ocr_missing = []
    try:
        import pytesseract
    except ImportError:
        ocr_missing.append("pytesseract")
    
    try:
        import pdf2image
    except ImportError:
        ocr_missing.append("pdf2image")
    
    if ocr_missing:
        console.print("[yellow]Missing OCR dependencies (needed for scanned PDFs):[/yellow]")
        for dep in ocr_missing:
            console.print(f"- {dep}")
        console.print("\nFor OCR support, install:")
        console.print(f"pip install {' '.join(ocr_missing)}")
        console.print("\nNote: You also need to install Tesseract OCR on your system:")
        console.print("https://github.com/UB-Mannheim/tesseract/wiki")
        console.print("For Korean documents, download the Korean language data pack")

    # Report missing dependencies
    if missing_deps:
        console.print("[red]Missing required dependencies:[/red]")
        for dep in missing_deps:
            console.print(f"- {dep}")
        console.print("\nPlease install them with:")
        console.print(f"pip install {' '.join(missing_deps)}")
        return False
    
    if optional_missing:
        console.print("[yellow]Missing optional dependencies:[/yellow]")
        for dep in optional_missing:
            console.print(f"- {dep}")
        console.print("\nFor full functionality, consider installing:")
        console.print(f"pip install {' '.join(optional_missing)}")
        
        if "spacy" in optional_missing:
            console.print("\nFor language processing, also install language models:")
            console.print("python -m spacy download en_core_web_sm")
            console.print("python -m spacy download ko_core_news_sm")
        
        console.print("\n[green]Required dependencies are satisfied, continuing...[/green]")
    else:
        console.print("[green]All dependencies are satisfied![/green]")
    
    return True


def main():
    """Main entry point."""
    try:
        # Check dependencies
        if not check_dependencies():
            sys.exit(1)
        
        # Show welcome message
        show_welcome_message()
        
        # Create and run app
        app = DocVectorizerApp()
        app.run()
        
    except Exception as e:
        console.print(f"[red]Fatal error: {str(e)}[/red]")
        logger.error(f"Fatal error: {str(e)}\n{traceback.format_exc()}")
        sys.exit(1)


if __name__ == "__main__":
    main()