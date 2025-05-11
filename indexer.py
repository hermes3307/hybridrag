#!/usr/bin/env python3
"""
Unified Hybrid RAG Document Indexer - Synchronous Version
All async processing has been removed for simplicity and stability
"""

import os
import sys
import json
import hashlib
import logging
from datetime import datetime
from pathlib import Path
from typing import Dict, Any, List, Optional, Literal
import argparse
import traceback
import time
import shutil
import re
from collections import deque

# Document processing
import PyPDF2
import docx
import pandas as pd
from pptx import Presentation

# ML/Vector libraries
from langchain_huggingface import HuggingFaceEmbeddings
from langchain_community.vectorstores import FAISS
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.schema import Document

# Graph database
from neo4j import GraphDatabase
from transformers import pipeline

# Configure logging
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(levelname)s - %(message)s',
    handlers=[
        logging.FileHandler('indexer.log', encoding='utf-8'),
        logging.StreamHandler()
    ]
)

# Define labeling strategies
LabelingStrategy = Literal["document_type", "domain_based", "hierarchical"]


class ProgressTracker:
    """Track progress with ETA calculation"""
    
    def __init__(self):
        self.start_time = time.time()
        self.file_times = deque(maxlen=50)
        self.current_file_start = None
        self.total_files = 0
        self.completed_files = 0
        self.current_phase = ""
        
    def start_tracking(self, total_files: int, phase: str = ""):
        """Start tracking progress"""
        self.total_files = total_files
        self.completed_files = 0
        self.start_time = time.time()
        self.current_phase = phase
        self.file_times.clear()
        
    def start_file(self):
        """Mark the start of file processing"""
        self.current_file_start = time.time()
    
    def complete_file(self, file_name: str = None):
        """Mark file as completed and update statistics"""
        if self.current_file_start:
            processing_time = time.time() - self.current_file_start
            self.file_times.append(processing_time)
            self.current_file_start = None
        
        self.completed_files += 1
        return self.get_progress_info(file_name)
    
    def get_progress_info(self, current_file: str = None) -> Dict[str, Any]:
        """Get current progress information"""
        if self.total_files == 0:
            return {}
        
        percentage = (self.completed_files / self.total_files) * 100
        elapsed_time = time.time() - self.start_time
        
        # Calculate ETA
        eta_seconds = None
        avg_time_per_file = None
        
        if self.file_times and self.completed_files > 0:
            avg_time_per_file = sum(self.file_times) / len(self.file_times)
            remaining_files = self.total_files - self.completed_files
            if remaining_files > 0:
                eta_seconds = avg_time_per_file * remaining_files
        
        # Format times
        elapsed_str = self._format_time(elapsed_time)
        eta_str = self._format_time(eta_seconds) if eta_seconds else "calculating..."
        avg_time_str = self._format_time(avg_time_per_file) if avg_time_per_file else "calculating..."
        
        return {
            'current_file': current_file,
            'completed': self.completed_files,
            'total': self.total_files,
            'percentage': percentage,
            'elapsed_time': elapsed_str,
            'eta': eta_str,
            'avg_time_per_file': avg_time_str,
            'remaining': self.total_files - self.completed_files,
            'phase': self.current_phase
        }
    
    def _format_time(self, seconds: float) -> str:
        """Format seconds into human-readable time"""
        if seconds is None:
            return "N/A"
        
        if seconds < 60:
            return f"{seconds:.1f}s"
        elif seconds < 3600:
            minutes = seconds / 60
            return f"{minutes:.1f}m"
        else:
            hours = seconds / 3600
            minutes = (seconds % 3600) / 60
            return f"{hours:.0f}h {minutes:.0f}m"
    
    def log_progress(self, logger, current_file: str = None, prefix: str = ""):
        """Log progress information"""
        info = self.get_progress_info(current_file)
        if not info:
            return
        
        progress_msg = (
            f"{prefix}[{info['completed']}/{info['total']}] "
            f"{info['percentage']:.1f}% | "
            f"Elapsed: {info['elapsed_time']} | "
            f"ETA: {info['eta']} | "
            f"Avg: {info['avg_time_per_file']}/file"
        )
        
        if current_file:
            progress_msg += f" | Current: {Path(current_file).name}"
        
        logger.info(progress_msg)


class DocumentLabeler:
    """Class to handle document labeling strategies"""
    
    def __init__(self, strategy: LabelingStrategy = "document_type"):
        self.strategy = strategy
        self.logger = logging.getLogger(__name__)
        
        # Define document type patterns
        self.type_patterns = {
            'report': ['report', '보고서', '리포트', 'analysis', '분석'],
            'contract': ['contract', '계약', 'agreement', '협약', '합의서'],
            'invoice': ['invoice', '청구서', '송장', 'receipt', '영수증'],
            'presentation': ['presentation', '발표', 'ppt', 'slides', '슬라이드'],
            'email': ['email', '이메일', 'mail', '메일'],
            'memo': ['memo', '메모', 'note', '공지', 'notice'],
            'resume': ['resume', '이력서', 'cv', '경력'],
            'manual': ['manual', '매뉴얼', 'guide', '가이드', 'instruction'],
            'research': ['research', '연구', 'paper', '논문', 'study'],
            'meeting': ['meeting', '회의', 'minutes', '회의록']
        }
        
        # Define domain patterns
        self.domain_patterns = {
            'finance': ['finance', '금융', 'bank', '은행', 'investment', '투자', 'budget', '예산'],
            'legal': ['legal', '법률', 'law', '법', 'regulation', '규정', 'compliance'],
            'technical': ['technical', '기술', 'engineering', '엔지니어링', 'development', '개발'],
            'medical': ['medical', '의료', 'health', '건강', 'patient', '환자', 'clinic'],
            'education': ['education', '교육', 'training', '훈련', 'course', '과정', 'learning'],
            'marketing': ['marketing', '마케팅', 'campaign', '캠페인', 'brand', '브랜드'],
            'hr': ['hr', '인사', 'employee', '직원', 'recruitment', '채용', 'performance'],
            'sales': ['sales', '영업', 'customer', '고객', 'deal', '거래', 'revenue'],
            'operations': ['operations', '운영', 'process', '프로세스', 'efficiency', '효율']
        }
        
        # Entity patterns for better extraction
        self.entity_patterns = {
            'organization': [r'\b[A-Z][A-Za-z\s&]+\s(?:Inc|Corp|Ltd|Company|회사|주식회사)\b'],
            'person': [r'\b[A-Z][a-z]+\s[A-Z][a-z]+\b', r'\b[가-힣]{2,4}\b'],
            'date': [r'\b\d{4}[-/]\d{1,2}[-/]\d{1,2}\b', r'\b\d{4}년\s?\d{1,2}월\s?\d{1,2}일\b'],
            'money': [r'\$[\d,]+\.?\d*', r'[\d,]+원', r'￦[\d,]+'],
            'percentage': [r'\d+\.?\d*%', r'\d+\.?\d*\s*퍼센트'],
            'email': [r'\b[A-Za-z0-9._%+-]+@[A-Za-z0-9.-]+\.[A-Z|a-z]{2,}\b'],
            'phone': [r'\b\d{3}[-.]?\d{3,4}[-.]?\d{4}\b', r'\b0\d{1,2}[-.]?\d{3,4}[-.]?\d{4}\b']
        }
    
    def determine_document_type(self, text: str, filename: str) -> str:
        """Determine document type based on content and filename"""
        text_lower = text.lower()
        filename_lower = filename.lower()
        
        # Check filename and content for type patterns
        for doc_type, patterns in self.type_patterns.items():
            for pattern in patterns:
                if pattern in filename_lower or pattern in text_lower[:1000]:
                    return doc_type
        
        # Check file extension
        ext = Path(filename).suffix.lower()
        if ext in ['.ppt', '.pptx']:
            return 'presentation'
        elif ext in ['.xlsx', '.xls']:
            return 'spreadsheet'
        elif ext in ['.pdf']:
            # Additional PDF analysis
            if 'invoice' in text_lower[:500] or '청구서' in text_lower[:500]:
                return 'invoice'
            elif 'contract' in text_lower[:500] or '계약' in text_lower[:500]:
                return 'contract'
        
        return 'document'  # Default type
    
    def determine_domain(self, text: str, filename: str) -> str:
        """Determine document domain based on content"""
        text_lower = text.lower()
        
        # Count domain keyword occurrences
        domain_scores = {}
        for domain, patterns in self.domain_patterns.items():
            score = 0
            for pattern in patterns:
                score += text_lower.count(pattern)
            domain_scores[domain] = score
        
        # Return domain with highest score, or 'general' if no clear domain
        if domain_scores:
            best_domain = max(domain_scores.items(), key=lambda x: x[1])
            if best_domain[1] > 0:
                return best_domain[0]
        
        return 'general'
    
    def extract_hierarchical_info(self, text: str, filename: str) -> Dict[str, str]:
        """Extract hierarchical information from document"""
        path_parts = Path(filename).parts
        
        hierarchy = {
            'department': 'unknown',
            'project': 'unknown',
            'category': 'unknown',
            'subcategory': 'unknown'
        }
        
        # Try to infer from path
        if len(path_parts) > 1:
            hierarchy['department'] = path_parts[0]
        if len(path_parts) > 2:
            hierarchy['project'] = path_parts[1]
        if len(path_parts) > 3:
            hierarchy['category'] = path_parts[2]
        
        return hierarchy
    
    def label_document(self, text: str, filename: str, metadata: Dict) -> Dict[str, Any]:
        """Apply labeling strategy to document"""
        result = {
            'labels': [],
            'properties': {},
            'relationships': []
        }
        
        if self.strategy == "document_type":
            doc_type = self.determine_document_type(text, filename)
            result['labels'] = ['Document', doc_type.capitalize()]
            result['properties'] = {
                'type': doc_type,
                'filename': filename,
                'created_at': datetime.now().isoformat()
            }
            result['relationships'].append({
                'type': 'IS_TYPE',
                'target': doc_type,
                'properties': {'confidence': 0.8}
            })
            
        elif self.strategy == "domain_based":
            domain = self.determine_domain(text, filename)
            doc_type = self.determine_document_type(text, filename)
            
            result['labels'] = ['Document', domain.capitalize()]
            result['properties'] = {
                'domain': domain,
                'type': doc_type,
                'filename': filename,
                'created_at': datetime.now().isoformat()
            }
            result['relationships'].append({
                'type': 'BELONGS_TO_DOMAIN',
                'target': domain,
                'properties': {'confidence': 0.8}
            })
            
        elif self.strategy == "hierarchical":
            hierarchy = self.extract_hierarchical_info(text, filename)
            doc_type = self.determine_document_type(text, filename)
            
            result['labels'] = ['Document', hierarchy['department'].capitalize()]
            result['properties'] = {
                'department': hierarchy['department'],
                'project': hierarchy['project'],
                'category': hierarchy['category'],
                'type': doc_type,
                'filename': filename,
                'created_at': datetime.now().isoformat()
            }
            
        return result
    
    def extract_enhanced_entities(self, text: str) -> List[Dict[str, Any]]:
        """Extract entities with enhanced patterns"""
        entities = []
        
        for entity_type, patterns in self.entity_patterns.items():
            for pattern in patterns:
                matches = re.finditer(pattern, text)
                for match in matches:
                    entity = {
                        'text': match.group(0),
                        'type': entity_type,
                        'start': match.start(),
                        'end': match.end(),
                        'confidence': 0.8
                    }
                    entities.append(entity)
        
        # Remove duplicates
        seen = set()
        unique_entities = []
        for entity in entities:
            key = f"{entity['text']}_{entity['type']}"
            if key not in seen:
                seen.add(key)
                unique_entities.append(entity)
        
        return unique_entities


class DocumentIndexer:
    """Unified Document Indexer for hybrid RAG system - Synchronous version"""
    
    def __init__(self, 
                 input_directory: str = "./batch_documents",
                 vector_db_path: str = "./vector_database",
                 neo4j_uri: str = "bolt://localhost:7687",
                 neo4j_user: str = "neo4j",
                 neo4j_password: str = "password",
                 chunk_size: int = 1000,
                 chunk_overlap: int = 200,
                 max_file_size_mb: int = 100,
                 labeling_strategy: Optional[LabelingStrategy] = None):
        
        # Initialize logger
        self.logger = logging.getLogger(__name__)
        
        # Configuration
        self.input_directory = Path(input_directory)
        self.vector_db_path = Path(vector_db_path)
        self.chunk_size = chunk_size
        self.chunk_overlap = chunk_overlap
        self.max_file_size_mb = max_file_size_mb
        self.labeling_strategy = labeling_strategy
        
        # Initialize document labeler if strategy provided
        self.labeler = DocumentLabeler(strategy=labeling_strategy) if labeling_strategy else None
        
        # Neo4j config
        self.neo4j_uri = neo4j_uri
        self.neo4j_user = neo4j_user
        self.neo4j_password = neo4j_password
        
        # Create output directory
        self.vector_db_path.mkdir(parents=True, exist_ok=True)
        
        # Initialize components
        self.embeddings = HuggingFaceEmbeddings(model_name="all-MiniLM-L6-v2")
        self.ner_pipeline = pipeline(
            "ner",
            model="dbmdz/bert-large-cased-finetuned-conll03-english",
            aggregation_strategy="simple"
        )
        
        # Neo4j connection
        self.driver = GraphDatabase.driver(neo4j_uri, auth=(neo4j_user, neo4j_password))
        
        # Text splitter
        self.text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=self.chunk_size,
            chunk_overlap=self.chunk_overlap,
            length_function=len,
        )
        
        # Tracking
        self.processed_files = {}
        self.processed_files_path = self.vector_db_path / "processed_files.json"
        self._load_processed_files()
        
        # Supported extensions
        self.supported_extensions = ['.pdf', '.docx', '.doc', '.txt', '.csv', '.xlsx', '.xls', '.pptx', '.ppt']
        
        # Track processing status
        self.processing_stats = {
            'success': 0,
            'failed': 0,
            'skipped': 0,
            'copied': 0
        }
        
        # Progress tracker
        self.progress_tracker = ProgressTracker()
    
    def _load_processed_files(self):
        """Load record of processed files"""
        if self.processed_files_path.exists():
            try:
                with open(self.processed_files_path, 'r', encoding='utf-8') as f:
                    self.processed_files = json.load(f)
            except (json.JSONDecodeError, Exception) as e:
                self.logger.error(f"Error loading processed files: {e}")
                self.processed_files = {}
                self._save_processed_files()
    
    def _save_processed_files(self):
        """Save record of processed files with error handling"""
        try:
            with open(self.processed_files_path, 'w', encoding='utf-8') as f:
                json.dump(self.processed_files, f, indent=2, ensure_ascii=False)
        except Exception as e:
            self.logger.error(f"Failed to save processed files: {e}")
    
    def _get_file_hash(self, file_path: Path) -> str:
        """Get MD5 hash of file"""
        try:
            hasher = hashlib.md5()
            with open(file_path, 'rb') as f:
                buf = f.read(65536)  # Read in 64kb chunks
                while len(buf) > 0:
                    hasher.update(buf)
                    buf = f.read(65536)
            return hasher.hexdigest()
        except Exception as e:
            self.logger.warning(f"Cannot hash {file_path.name}: {e}")
            return str(time.time())  # Return timestamp as fallback
    
    def _get_file_size_mb(self, file_path: Path) -> float:
        """Get file size in MB"""
        try:
            size_bytes = file_path.stat().st_size
            size_mb = size_bytes / (1024 * 1024)
            return size_mb
        except Exception as e:
            self.logger.warning(f"Cannot get file size for {file_path}: {e}")
            return 0.0
    
    def create_neo4j_document(self, file_path: str, metadata: Dict, text: str):
        """Create document node in Neo4j with optional labels"""
        with self.driver.session() as session:
            if self.labeler:
                # Enhanced version with labeling
                label_info = self.labeler.label_document(text, Path(file_path).name, metadata)
                labels = ':'.join(label_info['labels'])
                properties = {**metadata, **label_info['properties']}
                
                query = f"""
                MERGE (d:{labels} {{source: $source}})
                SET d += $properties
                """
                session.run(query, source=file_path, properties=properties)
                
                # Create relationships
                for rel in label_info['relationships']:
                    if rel['target'] != 'unknown':
                        target_query = f"""
                        MERGE (t:{rel['target'].capitalize()} {{name: $target_name}})
                        WITH t
                        MATCH (d:{labels} {{source: $source}})
                        MERGE (d)-[r:{rel['type']}]->(t)
                        SET r += $rel_properties
                        """
                        session.run(query=target_query, 
                                   source=file_path,
                                   target_name=rel['target'],
                                   rel_properties=rel['properties'])
            else:
                # Basic version without labeling
                query = """
                MERGE (d:Document {source: $source})
                SET d += $metadata
                """
                session.run(query, source=file_path, metadata=metadata)
    
    def extract_entities_and_relations(self, text: str) -> Dict[str, Any]:
        """Extract entities and relationships"""
        entities = []
        relations = []
        
        try:
            # Use enhanced entity extraction if available
            if self.labeler:
                enhanced_entities = self.labeler.extract_enhanced_entities(text)
                for ent in enhanced_entities:
                    entity = {
                        'id': hashlib.md5(f"{ent['text']}_{ent['type']}".encode()).hexdigest(),
                        'name': ent['text'],
                        'type': ent['type'],
                        'confidence': ent['confidence'],
                        'context': text[max(0, ent['start']-50):min(len(text), ent['end']+50)]
                    }
                    entities.append(entity)
            
            # Use HuggingFace NER for additional entities
            hf_entities = self.ner_pipeline(text[:1000])  # Limit for performance
            for ent in hf_entities:
                entity = {
                    'id': hashlib.md5(f"{ent['word']}_{ent['entity_group']}".encode()).hexdigest(),
                    'name': ent['word'],
                    'type': ent['entity_group'],
                    'confidence': ent['score'],
                    'context': text[max(0, ent['start']-50):min(len(text), ent['end']+50)]
                }
                entities.append(entity)
            
            # Extract relationships using patterns
            relation_patterns = [
                (r"(\w+)\s+(?:is|are)\s+(?:a|an)\s+(\w+)", "IS_A"),
                (r"(\w+)\s+(?:contains|includes)\s+(\w+)", "CONTAINS"),
                (r"(\w+)\s+(?:works at|employed by)\s+(\w+)", "WORKS_AT"),
                (r"(\w+)\s+(?:manages|leads)\s+(\w+)", "MANAGES"),
                (r"(\w+)\s+(?:reports to)\s+(\w+)", "REPORTS_TO"),
            ]
            
            for pattern, rel_type in relation_patterns:
                matches = re.finditer(pattern, text, re.IGNORECASE)
                for match in matches:
                    relation = {
                        'source_entity': match.group(1),
                        'target_entity': match.group(2),
                        'relation_type': rel_type,
                        'context': text[max(0, match.start()-50):min(len(text), match.end()+50)],
                        'confidence': 0.7
                    }
                    relations.append(relation)
            
        except Exception as e:
            self.logger.error(f"Entity extraction error: {e}")
        
        return {'entities': entities, 'relations': relations}
    
    def create_neo4j_entities(self, entities: List[Dict], source_file: str):
        """Create entities in Neo4j"""
        with self.driver.session() as session:
            for entity in entities:
                query = """
                MERGE (e:Entity {id: $id})
                SET e.name = $name,
                    e.type = $type,
                    e.source = $source,
                    e.context = $context,
                    e.confidence = $confidence,
                    e.created_at = datetime()
                """
                
                try:
                    session.run(query, 
                            id=entity['id'],
                            name=entity['name'],
                            type=entity['type'],
                            source=source_file,
                            context=entity['context'],
                            confidence=entity.get('confidence', 1.0))
                except Exception as e:
                    self.logger.debug(f"Could not create entity: {e}")
    
    def create_neo4j_relations(self, relations: List[Dict], source_file: str):
        """Create relationships in Neo4j"""
        with self.driver.session() as session:
            for relation in relations:
                try:
                    query = """
                    MATCH (e1:Entity {name: $source_entity})
                    MATCH (e2:Entity {name: $target_entity})
                    MERGE (e1)-[r:RELATES_TO {type: $relation_type}]->(e2)
                    SET r.context = $context,
                        r.source = $source,
                        r.created_at = datetime()
                    """
                    session.run(query,
                            source_entity=relation['source_entity'],
                            target_entity=relation['target_entity'],
                            relation_type=relation['relation_type'],
                            context=relation['context'],
                            source=source_file)
                except Exception as e:
                    self.logger.debug(f"Could not create relation: {e}")
    
    def extract_text(self, file_path: Path) -> str:
        """Extract text from document"""
        ext = file_path.suffix.lower()
        
        try:
            # Ensure file exists and is readable
            if not file_path.exists():
                self.logger.error(f"File not found: {file_path}")
                return ""
            
            if ext == '.pdf':
                return self._extract_text_from_pdf(file_path)
            elif ext == '.docx':
                return self._extract_text_from_docx(file_path)
            elif ext == '.txt':
                return self._extract_text_from_txt(file_path)
            elif ext in ['.xlsx', '.xls']:
                return self._extract_text_from_excel(file_path)
            elif ext in ['.pptx', '.ppt']:
                return self._extract_text_from_pptx(file_path)
            elif ext == '.csv':
                return self._extract_text_from_csv(file_path)
            else:
                self.logger.warning(f"Unsupported file type: {ext}")
                return ""
        except Exception as e:
            self.logger.error(f"Error extracting text from {file_path.name}: {e}")
            return ""
    
    def _extract_text_from_pdf(self, file_path: Path) -> str:
        """Extract text from PDF"""
        text = ""
        try:
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                for page in pdf_reader.pages:
                    page_text = page.extract_text()
                    if page_text:
                        text += page_text + "\n"
        except Exception as e:
            self.logger.error(f"PDF extraction error: {e}")
        return text
    
    def _extract_text_from_docx(self, file_path: Path) -> str:
        """Extract text from DOCX"""
        text = ""
        try:
            doc = docx.Document(str(file_path))
            for paragraph in doc.paragraphs:
                text += paragraph.text + "\n"
        except Exception as e:
            self.logger.error(f"DOCX extraction error for {file_path.name}: {e}")
        return text
    
    def _extract_text_from_txt(self, file_path: Path) -> str:
        """Extract text from TXT"""
        encodings = ['utf-8', 'cp949', 'cp1252', 'latin1']
        for encoding in encodings:
            try:
                with open(file_path, 'r', encoding=encoding) as f:
                    text = f.read()
                    if text.strip():
                        return text
            except UnicodeDecodeError:
                continue
        return ""
    
    def _extract_text_from_csv(self, file_path: Path) -> str:
        """Extract text from CSV with multiple encoding attempts"""
        encodings = ['utf-8', 'cp949', 'cp1252', 'latin1', 'euc-kr', 'iso-8859-1', 'utf-16', 'utf-32']
        
        # First check file size
        file_size = file_path.stat().st_size
        self.logger.debug(f"CSV file {file_path.name} size: {file_size} bytes")
        
        if file_size == 0:
            self.logger.warning(f"CSV file {file_path.name} is empty (0 bytes)")
            return ""
        
        for encoding in encodings:
            try:
                df = pd.read_csv(file_path, encoding=encoding)
                # Only proceed if we have data
                if not df.empty:
                    self.logger.info(f"Successfully read CSV {file_path.name} with {encoding} encoding")
                    return df.to_string()
            except UnicodeDecodeError as e:
                self.logger.debug(f"Unicode error with {encoding} for {file_path.name}: {e}")
                continue
            except Exception as e:
                self.logger.debug(f"CSV read error with {encoding} for {file_path.name}: {e}")
                continue
        
        # Try with 'errors' parameter as fallback
        try:
            df = pd.read_csv(file_path, encoding='utf-8', errors='replace')
            self.logger.info(f"Read CSV {file_path.name} with error replacement")
            return df.to_string()
        except Exception as e:
            self.logger.error(f"CSV extraction failed for {file_path.name}: {e}")
        
        # Last resort: try to read as bytes and decode
        try:
            with open(file_path, 'rb') as f:
                content = f.read()
                # Try to decode with errors='replace'
                text = content.decode('utf-8', errors='replace')
                self.logger.info(f"Read CSV {file_path.name} as raw text")
                return text
        except Exception as e:
            self.logger.error(f"Final CSV extraction attempt failed for {file_path.name}: {e}")
        
        return ""
    
    def _extract_text_from_excel(self, file_path: Path) -> str:
        """Extract text from Excel with better error handling"""
        text = ""
        try:
            # Try different engines for different Excel formats
            try:
                df = pd.read_excel(file_path, sheet_name=None)
            except Exception:
                # Try with openpyxl engine explicitly
                df = pd.read_excel(file_path, sheet_name=None, engine='openpyxl')
            
            for sheet_name, sheet_df in df.items():
                if not sheet_df.empty:
                    text += f"Sheet: {sheet_name}\n"
                    text += sheet_df.to_string() + "\n\n"
            
            if not text.strip():
                self.logger.warning(f"Excel file {file_path.name} appears to be empty")
                
        except Exception as e:
            self.logger.error(f"Excel extraction error for {file_path.name}: {e}")
        return text
    
    def _extract_text_from_pptx(self, file_path: Path) -> str:
        """Extract text from PowerPoint"""
        text = ""
        try:
            prs = Presentation(str(file_path))
            for slide_num, slide in enumerate(prs.slides):
                text += f"[Slide {slide_num + 1}]\n"
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
        except Exception as e:
            self.logger.error(f"PowerPoint extraction error: {e}")
        return text
    
    def process_file(self, file_path: Path) -> bool:
        """Process a single file - synchronous processing"""
        try:
            start_time = time.time()
            file_str = str(file_path)
            
            # Check if file exists
            if not file_path.exists():
                self.logger.warning(f"File doesn't exist: {file_path}")
                self.processing_stats['skipped'] += 1
                return False
            
            # Check file size
            file_size = file_path.stat().st_size
            size_mb = self._get_file_size_mb(file_path)
            
            # Debug: Log actual file size
            self.logger.debug(f"File {file_path.name}: {file_size} bytes ({size_mb:.3f} MB)")
            
            if file_size == 0:
                self.logger.warning(f"Skipping empty file: {file_path.name}")
                self.processing_stats['skipped'] += 1
                return False
            
            if size_mb > self.max_file_size_mb:
                self.logger.warning(f"Skipping large file: {file_path.name} ({size_mb:.3f} MB)")
                self.processing_stats['skipped'] += 1
                return False
            
            # Get file hash
            file_hash = self._get_file_hash(file_path)
            
            # Check if already processed
            if file_str in self.processed_files:
                if self.processed_files[file_str]['hash'] == file_hash:
                    self.logger.info(f"Already processed: {file_path.name}")
                    return True
            
            self.logger.info(f"Processing: {file_path.name} ({size_mb:.3f} MB)")
            
            # Extract text
            text = self.extract_text(file_path)
            if not text or not text.strip():
                self.logger.warning(f"No text extracted from {file_path.name}")
                self.processing_stats['failed'] += 1
                return False
            
            # Split into chunks
            chunks = self.text_splitter.split_text(text)
            if not chunks:
                self.logger.warning(f"No chunks created for {file_path.name}")
                self.processing_stats['failed'] += 1
                return False
            
            # Create document metadata
            metadata = {
                'source': file_str,
                'filename': file_path.name,
                'file_hash': file_hash,
                'processed_date': datetime.now().isoformat(),
                'chunks': len(chunks)
            }
            
            # Create document node in Neo4j
            try:
                self.create_neo4j_document(file_str, metadata, text)
            except Exception as e:
                self.logger.error(f"Neo4j document creation failed for {file_path.name}: {e}")
            
            # Process chunks for vector store
            documents = []
            for i, chunk in enumerate(chunks):
                try:
                    doc_metadata = metadata.copy()
                    doc_metadata['chunk_index'] = i
                    documents.append(Document(page_content=chunk, metadata=doc_metadata))
                    
                    # Extract entities and relations
                    kg_data = self.extract_entities_and_relations(chunk)
                    
                    # Store in Neo4j
                    if kg_data.get('entities'):
                        self.create_neo4j_entities(kg_data['entities'], file_str)
                    if kg_data.get('relations'):
                        self.create_neo4j_relations(kg_data['relations'], file_str)
                except Exception as e:
                    self.logger.error(f"Chunk processing failed for {file_path.name}, chunk {i}: {e}")
            
            # Add to vector store
            if documents:
                try:
                    if not hasattr(self, 'vector_store'):
                        self.vector_store = FAISS.from_documents(documents, self.embeddings)
                    else:
                        self.vector_store.add_documents(documents)
                    
                    # Save immediately
                    self._save_vector_store()
                except Exception as e:
                    self.logger.error(f"Vector store update failed for {file_path.name}: {e}")
            
            # Update processed files
            self.processed_files[file_str] = {
                'hash': file_hash,
                'processed_date': datetime.now().isoformat(),
                'chunks': len(chunks),
                'processing_time': time.time() - start_time
            }
            self.processing_stats['success'] += 1
            self._save_processed_files()
            
            processing_time = time.time() - start_time
            self.logger.info(f"Successfully processed {file_path.name}: {len(chunks)} chunks in {processing_time:.2f}s")
            return True
            
        except Exception as e:
            self.logger.error(f"Error processing {file_path}: {e}")
            self.logger.error(traceback.format_exc())
            self.processing_stats['failed'] += 1
            return False
    
    def scan_and_copy_directory(self, source_directory: str, target_directory: str = "./batch_documents") -> List[Path]:
        """Scan source directory recursively and copy files to target directory"""
        source_path = Path(source_directory)
        target_path = Path(target_directory)
        
        if not source_path.exists():
            self.logger.error(f"Source directory not found: {source_directory}")
            return []
        
        # Create target directory if it doesn't exist
        target_path.mkdir(parents=True, exist_ok=True)
        
        self.logger.info(f"Scanning directory: {source_directory}")
        print(f"Scanning directory structure...")
        
        # Collect all files to process
        files_to_copy = []
        exclude_patterns = ['.tmp', '.temp', '~$', '.lock', 'desktop.ini', 'thumbs.db', '.DS_Store']
        
        for ext in self.supported_extensions:
            for file_path in source_path.rglob(f"*{ext}"):
                # Skip system/temp files
                if any(pattern in file_path.name.lower() for pattern in exclude_patterns):
                    continue
                
                # Skip very small files
                try:
                    if file_path.stat().st_size < 10:
                        continue
                except:
                    continue
                
                files_to_copy.append(file_path)
        
        total_files = len(files_to_copy)
        if total_files == 0:
            self.logger.warning("No supported files found in source directory")
            return []
        
        self.logger.info(f"Found {total_files} files to copy")
        
        # Initialize progress tracker
        self.progress_tracker.start_tracking(total_files, "copying")
        
        copied_files = []
        
        for idx, file_path in enumerate(files_to_copy, 1):
            try:
                self.progress_tracker.start_file()
                
                # Generate target path
                relative_path = file_path.relative_to(source_path)
                target_file_path = target_path / relative_path
                
                # Create subdirectories if needed
                target_file_path.parent.mkdir(parents=True, exist_ok=True)
                
                # Copy file
                shutil.copy2(file_path, target_file_path)
                copied_files.append(target_file_path)
                
                self.processing_stats['copied'] += 1
                
                # Update progress
                self.progress_tracker.complete_file(str(file_path))
                
                # Log progress periodically
                if idx % 10 == 0 or idx == total_files:
                    self.progress_tracker.log_progress(self.logger, str(file_path), "Copying ")
                
            except Exception as e:
                self.logger.error(f"Error copying {file_path}: {e}")
                self.progress_tracker.complete_file(str(file_path))
        
        final_info = self.progress_tracker.get_progress_info()
        self.logger.info(
            f"Copying completed: {final_info['completed']}/{final_info['total']} files "
            f"in {final_info['elapsed_time']}"
        )
        
        return copied_files
    
    def index_directory(self, batch_size: int = 10):
        """Index all documents in the directory - synchronous processing"""
        if not self.input_directory.exists():
            self.logger.error(f"Directory not found: {self.input_directory}")
            return
        
        # Find all supported files
        files_to_process = []
        for ext in self.supported_extensions:
            files_to_process.extend(self.input_directory.rglob(f"*{ext}"))
        
        # Filter and sort files
        valid_files = []
        for file_path in files_to_process:
            try:
                if file_path.exists() and file_path.is_file():
                    valid_files.append(file_path)
            except Exception as e:
                self.logger.warning(f"Cannot access file {file_path}: {e}")
                continue
        
        # Sort files by size (process smaller files first)
        valid_files.sort(key=lambda f: self._get_file_size_mb(f) if f.exists() else 0)
        
        total_files = len(valid_files)
        self.logger.info(f"Found {total_files} files to process")
        
        if total_files == 0:
            self.logger.warning("No files found to index")
            return
        
        # Initialize progress tracker
        self.progress_tracker.start_tracking(total_files, "indexing")
        
        # Process each file
        for i, file_path in enumerate(valid_files, 1):
            try:
                self.progress_tracker.start_file()
                
                # Process file
                success = self.process_file(file_path)
                
                # Update progress
                self.progress_tracker.complete_file(str(file_path))
                
                # Log progress periodically
                if i % 5 == 0 or i == total_files:
                    info = self.progress_tracker.get_progress_info()
                    self.logger.info(
                        f"Progress: {info['completed']}/{info['total']} files "
                        f"({info['percentage']:.1f}%) | "
                        f"Success: {self.processing_stats['success']}, "
                        f"Failed: {self.processing_stats['failed']}, "
                        f"Skipped: {self.processing_stats['skipped']}"
                    )
                
                # Save progress periodically
                if i % batch_size == 0:
                    try:
                        self._save_vector_store()
                        self._save_processed_files()
                        self.logger.info(f"Batch saved at {i} files")
                    except Exception as e:
                        self.logger.error(f"Error saving batch: {e}")
            
            except Exception as e:
                self.logger.error(f"Error processing file {i}/{total_files} ({file_path}): {e}")
                self.progress_tracker.complete_file(str(file_path))
                continue
        
        # Final save
        try:
            self._save_vector_store()
            self._save_processed_files()
        except Exception as e:
            self.logger.error(f"Error in final save: {e}")
        
        final_info = self.progress_tracker.get_progress_info()
        self.logger.info(
            f"Indexing completed: {final_info['completed']}/{final_info['total']} files "
            f"in {final_info['elapsed_time']}"
        )
        self._log_stats()
    
    def scan_copy_and_index(self, source_directory: str, target_directory: str = "./batch_documents", 
                           batch_size: int = 10):
        """Scan source directory, copy files to target, and index"""
        self.logger.info(f"Starting scan-copy-index from {source_directory} to {target_directory}")
        if self.labeling_strategy:
            self.logger.info(f"Using labeling strategy: {self.labeling_strategy}")
        
        # Reset stats
        self.processing_stats = {
            'success': 0,
            'failed': 0,
            'skipped': 0,
            'copied': 0
        }
        
        # Scan and copy files
        copied_files = self.scan_and_copy_directory(source_directory, target_directory)
        
        if not copied_files:
            self.logger.warning("No files found to process")
            return
        
        # Update input directory and index
        original_input_dir = self.input_directory
        self.input_directory = Path(target_directory)
        
        self.index_directory(batch_size)
        
        # Restore original input directory
        self.input_directory = original_input_dir
        
        self.logger.info("Scan-copy-index process completed")
        self._log_stats()
    
    def get_statistics(self) -> Dict:
        """Get indexing statistics"""
        stats = {
            'total_processed': len(self.processed_files),
            'vector_db_path': str(self.vector_db_path),
            'input_directory': str(self.input_directory),
            'processing_stats': self.processing_stats.copy(),
        }
        
        if self.labeling_strategy:
            stats['labeling_strategy'] = self.labeling_strategy
        
        # Get Neo4j statistics
        with self.driver.session() as session:
            try:
                # Document statistics
                if self.labeling_strategy:
                    doc_stats_query = """
                    MATCH (d:Document)
                    RETURN labels(d) as labels, count(d) as count
                    """
                    result = session.run(doc_stats_query)
                    stats['documents_by_label'] = {
                        str(record['labels']): record['count'] 
                        for record in result
                    }
                
                # Overall counts
                result = session.run("MATCH (d:Document) RETURN count(d) as count")
                stats['total_documents'] = result.single()['count']
                
                result = session.run("MATCH (e:Entity) RETURN count(e) as count")
                stats['total_entities'] = result.single()['count']
                
                result = session.run("MATCH ()-[r:RELATES_TO]->() RETURN count(r) as count")
                stats['total_relations'] = result.single()['count']
            except Exception as e:
                self.logger.error(f"Error getting Neo4j statistics: {e}")
                stats['total_documents'] = 0
                stats['total_entities'] = 0
                stats['total_relations'] = 0
        
        return stats
    
    def _log_stats(self):
        """Log current processing statistics"""
        self.logger.info(f"Processing stats - Success: {self.processing_stats['success']}, "
                        f"Failed: {self.processing_stats['failed']}, Skipped: {self.processing_stats['skipped']}, "
                        f"Copied: {self.processing_stats['copied']}")
    
    def _save_vector_store(self):
        """Save vector store to disk with error handling"""
        try:
            if hasattr(self, 'vector_store'):
                self.vector_store.save_local(str(self.vector_db_path))
                self.logger.debug("Vector store saved")
        except Exception as e:
            self.logger.error(f"Failed to save vector store: {e}")
    
    def close(self):
        """Close connections and cleanup"""
        self.driver.close()


def main():
    """Main function with interactive or command-line interface"""
    parser = argparse.ArgumentParser(description="Unified Hybrid RAG Document Indexer - Synchronous Version")
    parser.add_argument("--interactive", action="store_true", help="Run in interactive mode")
    parser.add_argument("--scan-dir", help="Source directory to scan and copy from")
    parser.add_argument("--input-dir", default="./batch_documents", help="Input directory containing documents")
    parser.add_argument("--target-dir", default="./batch_documents", help="Target directory for copied files")
    parser.add_argument("--vector-db", default="./vector_database", help="Vector database output path")
    parser.add_argument("--neo4j-uri", default="bolt://localhost:7687", help="Neo4j URI")
    parser.add_argument("--neo4j-user", default="neo4j", help="Neo4j username")
    parser.add_argument("--neo4j-password", default="password", help="Neo4j password")
    parser.add_argument("--batch-size", type=int, default=10, help="Batch size for saving progress")
    parser.add_argument("--labeling-strategy", choices=["document_type", "domain_based", "hierarchical"], 
                       help="Labeling strategy for graph database")
    
    args = parser.parse_args()
    
    # Interactive mode
    if args.interactive or len(sys.argv) == 1:
        print("=== Unified Hybrid RAG Document Indexer (Synchronous Version) ===")
        
        # Get configuration
        neo4j_password = input("Enter Neo4j password [password]: ").strip() or "password"
        
        # Get labeling strategy
        print("\n=== Document Labeling Strategy (Optional) ===")
        print("1. Skip labeling (basic indexing)")
        print("2. Document Type Based (Report, Contract, Invoice, etc.)")
        print("3. Domain Based (Finance, Legal, Technical, etc.)")
        print("4. Hierarchical (Department/Project/Category)")
        
        choice = input("\nSelect strategy (1-4) [1]: ").strip() or "1"
        
        labeling_strategy = None
        if choice == "2":
            labeling_strategy = "document_type"
        elif choice == "3":
            labeling_strategy = "domain_based"
        elif choice == "4":
            labeling_strategy = "hierarchical"
        
        # Initialize indexer
        indexer = DocumentIndexer(
            neo4j_password=neo4j_password,
            labeling_strategy=labeling_strategy
        )
        
        if labeling_strategy:
            print(f"\nIndexer initialized with {labeling_strategy} labeling strategy")
        else:
            print("\nIndexer initialized (basic mode)")
        
        # Get source directory
        source_dir = input("\nEnter source directory path (or press Enter to skip): ").strip()
        
        if source_dir:
            # Scan, copy, and index
            try:
                indexer.scan_copy_and_index(source_directory=source_dir)
                print("\nIndexing completed successfully!")
            except Exception as e:
                print(f"Error during indexing: {e}")
        else:
            # Just index existing files
            try:
                indexer.index_directory()
                print("\nIndexing completed successfully!")
            except Exception as e:
                print(f"Error during indexing: {e}")
        
        # Display statistics
        stats = indexer.get_statistics()
        print("\n=== Indexing Statistics ===")
        print(f"Total documents: {stats.get('total_documents', 0)}")
        print(f"Total entities: {stats.get('total_entities', 0)}")
        print(f"Total relationships: {stats.get('total_relations', 0)}")
        
        if labeling_strategy and 'documents_by_label' in stats:
            print("\n=== Documents by Label ===")
            for labels, count in stats.get('documents_by_label', {}).items():
                print(f"  {labels}: {count}")
        
        proc_stats = stats['processing_stats']
        print("\n=== Processing Summary ===")
        print(f"Successfully processed: {proc_stats['success']}")
        print(f"Failed: {proc_stats['failed']}")
        print(f"Skipped: {proc_stats['skipped']}")
        if source_dir:
            print(f"Files copied: {proc_stats['copied']}")
        
        indexer.close()
        
    else:
        # Command-line mode
        indexer = DocumentIndexer(
            input_directory=args.input_dir,
            vector_db_path=args.vector_db,
            neo4j_uri=args.neo4j_uri,
            neo4j_user=args.neo4j_user,
            neo4j_password=args.neo4j_password,
            labeling_strategy=args.labeling_strategy
        )
        
        try:
            if args.scan_dir:
                # Scan, copy, and index
                print(f"Scanning directory: {args.scan_dir}")
                if args.labeling_strategy:
                    print(f"Using labeling strategy: {args.labeling_strategy}")
                
                indexer.scan_copy_and_index(
                    source_directory=args.scan_dir,
                    target_directory=args.target_dir,
                    batch_size=args.batch_size
                )
            else:
                # Standard indexing
                print(f"Indexing directory: {args.input_dir}")
                if args.labeling_strategy:
                    print(f"Using labeling strategy: {args.labeling_strategy}")
                
                indexer.index_directory(batch_size=args.batch_size)
            
            # Print statistics
            stats = indexer.get_statistics()
            print("\nIndexing completed!")
            print(f"Documents: {stats.get('total_documents', 0)}")
            print(f"Entities: {stats.get('total_entities', 0)}")
            print(f"Relationships: {stats.get('total_relations', 0)}")
            
            proc_stats = stats['processing_stats']
            print(f"\nProcessed: {proc_stats['success']} success, {proc_stats['failed']} failed")
            
        except Exception as e:
            print(f"Error: {e}")
            traceback.print_exc()
        finally:
            indexer.close()


if __name__ == "__main__":
    main()